"""Microsoft Office 文档提取器：DOCX 与 PPTX。

DOCX/PPTX 优先使用 lxml (libxml2 C 扩展) 直接解析 OOXML XML，
绕开 python-docx/python-pptx 的对象封装，性能提升 5-10x。
lxml 不可用时回退到 python-docx/python-pptx。

详见 :mod:`fuscan.extractors._ooxml_xml`。
"""

from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Any

from typing_extensions import override

from fuscan.extractors._ooxml_xml import extract_docx_text, extract_pptx_text
from fuscan.extractors.base import Extractor, ExtractorError, SpeedTier

__all__ = ["DocxExtractor", "PptxExtractor"]

logger = logging.getLogger(__name__)


def _lxml_available() -> bool:
    """检查 lxml 是否可导入。"""
    try:
        import lxml  # noqa: F401
    except ImportError:
        return False
    return True


class DocxExtractor(Extractor):
    """DOCX 文档文本提取器。

    优先使用 lxml 直接解析 ``word/document.xml``，
    回退到 python-docx（功能等价但较慢）。
    """

    @property
    @override
    def supported_extensions(self) -> tuple[str, ...]:
        """返回 DOCX 提取器支持的扩展名。"""
        return ("docx",)

    @property
    @override
    def speed_tier(self) -> SpeedTier:
        """lxml 直接解析 XML 为 T2 快速；回退 python-docx 为 T3 中速。"""
        return SpeedTier.FAST if _lxml_available() else SpeedTier.MEDIUM

    @override
    @property
    def display_name(self) -> str:
        """返回提取器的中文显示名称。"""
        return "Word（DOCX）"

    @override
    @property
    def engine_info(self) -> str:
        """iter-139：lxml 可用时优先使用，回退 python-docx。"""
        return "lxml" if _lxml_available() else "python-docx"

    @override
    def extract(self, path: Path) -> str:
        """提取 DOCX 段落、表格与页眉页脚文本。"""
        try:
            data = path.read_bytes()
        except OSError as exc:
            raise ExtractorError(f"文件读取失败: {path}: {exc}") from exc
        return self.extract_from_bytes(data)

    @override
    def extract_from_bytes(self, data: bytes) -> str:
        """从内存字节提取 DOCX 文本。"""
        if _lxml_available():
            try:
                return extract_docx_text(data)
            except Exception as exc:
                # ZIP 损坏或 XML 严重损坏，回退到 python-docx 再试一次
                if _is_zip_error(exc):
                    raise ExtractorError(f"DOCX 解析失败: {exc}") from exc
                logger.debug("lxml 解析 DOCX 失败，回退 python-docx: %s", exc)

        try:
            from docx import Document
        except ImportError as exc:
            raise ExtractorError("python-docx 未安装，无法提取 DOCX") from exc

        try:
            doc = Document(io.BytesIO(data))
        except Exception as exc:
            raise ExtractorError(f"DOCX 解析失败: {exc}") from exc

        parts: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
        for table in doc.tables:
            for row in table.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    parts.append("\t".join(row_texts))
        parts.extend(self._extract_docx_sections(doc.sections))
        return "\n".join(parts)

    @staticmethod
    def _extract_docx_sections(sections: Any) -> list[str]:
        """从 DOCX 节的页眉页脚提取文本。"""
        texts: list[str] = []
        for section in sections:
            for header_footer in (section.header, section.footer):
                for para in header_footer.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        return texts


class PptxExtractor(Extractor):
    """PPTX 演示文稿文本提取器。

    优先使用 lxml 直接解析 ``ppt/slides/slideN.xml``，
    回退到 python-pptx（功能等价但较慢）。
    """

    @property
    @override
    def supported_extensions(self) -> tuple[str, ...]:
        """返回 PPTX 提取器支持的扩展名。"""
        return ("pptx",)

    @property
    @override
    def speed_tier(self) -> SpeedTier:
        """lxml 直接解析 XML 为 T2 快速；回退 python-pptx 为 T4 慢速。"""
        return SpeedTier.FAST if _lxml_available() else SpeedTier.SLOW

    @override
    @property
    def display_name(self) -> str:
        """返回提取器的中文显示名称。"""
        return "PowerPoint（PPTX）"

    @override
    @property
    def engine_info(self) -> str:
        """iter-139：lxml 可用时优先使用，回退 python-pptx。"""
        return "lxml" if _lxml_available() else "python-pptx"

    @override
    def extract(self, path: Path) -> str:
        """提取 PPTX 幻灯片文本框、表格与备注文本。"""
        try:
            data = path.read_bytes()
        except OSError as exc:
            raise ExtractorError(f"文件读取失败: {path}: {exc}") from exc
        return self.extract_from_bytes(data)

    @override
    def extract_from_bytes(self, data: bytes) -> str:
        """从内存字节提取 PPTX 文本。"""
        if _lxml_available():
            try:
                return extract_pptx_text(data)
            except Exception as exc:
                if _is_zip_error(exc):
                    raise ExtractorError(f"PPTX 解析失败: {exc}") from exc
                logger.debug("lxml 解析 PPTX 失败，回退 python-pptx: %s", exc)

        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ExtractorError("python-pptx 未安装，无法提取 PPTX") from exc

        try:
            prs = Presentation(io.BytesIO(data))
        except Exception as exc:
            raise ExtractorError(f"PPTX 解析失败: {exc}") from exc

        parts: list[str] = []
        for slide_index, slide in enumerate(prs.slides, 1):
            slide_texts = self._extract_slide(slide)
            if slide_texts:
                parts.append(f"--- 幻灯片 {slide_index} ---")
                parts.extend(slide_texts)
        return "\n".join(parts)

    def _extract_slide(self, slide: object) -> list[str]:
        """提取单张幻灯片的文本。"""
        texts: list[str] = []
        for shape in slide.shapes:  # pyrefly: ignore [missing-attribute]
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        texts.append("\t".join(row_texts))
        if slide.has_notes_slide:  # pyrefly: ignore [missing-attribute]
            notes_text = slide.notes_slide.notes_text_frame.text.strip()  # pyrefly: ignore [missing-attribute]
            if notes_text:
                texts.append(f"[备注] {notes_text}")
        return texts


def _is_zip_error(exc: Exception) -> bool:
    """判断异常是否为 ZIP 格式错误（无法恢复，不应回退）。"""
    import zipfile

    if isinstance(exc, zipfile.BadZipFile):
        return True
    # lxml 的 XMLSyntaxError 在 recover=True 下通常不抛出，但严重损坏时可能抛出
    return isinstance(exc, type(None))  # 其他异常允许回退
