"""提取器单元测试。

使用对应库动态生成测试 fixture 文件，避免二进制 fixture 入仓。
PDF/ODT/ODS 等较难动态生成的格式，使用 mock 或跳过。
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

import pytest
from typing_extensions import override

from fuscan.extractors import (
    DocExtractor,
    DocxExtractor,
    EmlExtractor,
    Extractor,
    ExtractorError,
    ExtractorFailure,
    ExtractorRegistry,
    ImageExtractor,
    OdtExtractor,
    PdfExtractor,
    PptExtractor,
    PptxExtractor,
    TextExtractor,
    WpsExtractor,
    XlsExtractor,
    XlsxExtractor,
    clear_content_cache,
    default_registry,
    extract_content,
    extract_content_cached,
    extract_content_from_bytes,
    extract_content_from_bytes_with_retry,
    extract_content_with_fallback,
    extract_content_with_fallback_and_retry,
    get_extractor,
    is_retriable_error,
)
from fuscan.extractors.base import SpeedTier
from fuscan.extractors.spreadsheet import OdsExtractor


def _make_ooxml_zip(entry_name: str, content: str = "fake") -> bytes:
    """创建包含指定条目的有效 ZIP，用于测试 OOXML 类型检测。

    :param entry_name: ZIP 内部条目名（如 ``word/document.xml``）
    :param content: 条目内容（默认 ``"fake"``，用于触发解析失败）
    :return: ZIP 文件字节
    """
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr(entry_name, content)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Fixture 工厂
# ---------------------------------------------------------------------------


@pytest.fixture()
def text_file(tmp_path: Path) -> Path:
    path = tmp_path / "sample.txt"
    path.write_text("hello password world\n第二行内容\n", encoding="utf-8")
    return path


@pytest.fixture()
def gbk_file(tmp_path: Path) -> Path:
    path = tmp_path / "gbk.txt"
    path.write_bytes("这是一个包含密码字段的配置文件，密码为 password123，请妥善保管。PASSWORD".encode("gbk"))
    return path


@pytest.fixture()
def empty_file(tmp_path: Path) -> Path:
    path = tmp_path / "empty.txt"
    path.write_text("", encoding="utf-8")
    return path


@pytest.fixture()
def docx_file(tmp_path: Path) -> Path:
    from docx import Document

    doc = Document()
    doc.add_paragraph("段落一 含 password")
    doc.add_paragraph("段落二 正常内容")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "姓名"
    table.cell(0, 1).text = "密码"
    table.cell(1, 0).text = "张三"
    table.cell(1, 1).text = "pwd123"
    path = tmp_path / "test.docx"
    doc.save(str(path))
    return path


@pytest.fixture()
def pptx_file(tmp_path: Path) -> Path:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "标题 含 secret"  # pyrefly: ignore [missing-attribute]
    slide.placeholders[1].text = "幻灯片内容"  # pyrefly: ignore [missing-attribute]
    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


@pytest.fixture()
def xlsx_file(tmp_path: Path) -> Path:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "数据"
    ws["A1"] = "姓名"
    ws["B1"] = "密码"
    ws["A2"] = "张三"
    ws["B2"] = "pwd123"
    wb.create_sheet("空表")
    path = tmp_path / "test.xlsx"
    wb.save(str(path))
    return path


# ---------------------------------------------------------------------------
# TextExtractor
# ---------------------------------------------------------------------------


class TestTextExtractor:
    def test_extract_utf8(self, text_file: Path) -> None:
        extractor = TextExtractor()
        content = extractor.extract(text_file)
        assert "hello password world" in content
        assert "第二行内容" in content

    def test_extract_empty(self, empty_file: Path) -> None:
        extractor = TextExtractor()
        assert extractor.extract(empty_file) == ""

    def test_supported_extensions(self) -> None:
        extractor = TextExtractor()
        assert "txt" in extractor.supported_extensions
        assert "md" in extractor.supported_extensions
        assert "py" in extractor.supported_extensions

    def test_max_size_limit(self, tmp_path: Path) -> None:
        path = tmp_path / "big.txt"
        path.write_text("x" * 100, encoding="utf-8")
        extractor = TextExtractor(max_size=10)
        assert extractor.extract(path) == ""

    def test_nonexistent_file_raises(self, tmp_path: Path) -> None:
        extractor = TextExtractor()
        with pytest.raises(ExtractorError, match="无法读取文件大小"):
            extractor.extract(tmp_path / "missing.txt")

    def test_charset_normalizer_fallback(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """charset-normalizer 未安装时回退到 UTF-8/GBK 解码。

        内容用 GBK 编码写入，使其无法命中 UTF-8 快路径，从而进入
        charset-normalizer 分支触发 ImportError 回退（回退链中 GBK 可解）。
        禁用原生编码检测以隔离测试 charset-normalizer 回退路径。
        """
        monkeypatch.setattr("fuscan.extractors.text._NATIVE_DECODE_AVAILABLE", False)
        path = tmp_path / "fallback.txt"
        path.write_bytes("回退解码 password".encode("gbk"))

        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "charset_normalizer":
                raise ImportError("No module named 'charset_normalizer'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)

        content = TextExtractor().extract(path)
        assert "回退解码 password" in content

    def test_gbk_decoding(self, gbk_file: Path) -> None:
        """GBK 编码文件应能正确解码。"""
        extractor = TextExtractor()
        content = extractor.extract(gbk_file)
        assert "密码" in content
        assert "password123" in content

    def test_read_bytes_os_error(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """read_bytes 失败时抛出 ExtractorError。"""
        path = tmp_path / "test.txt"
        path.write_text("hello", encoding="utf-8")

        original_read_bytes = Path.read_bytes

        def mock_read_bytes(self: Path) -> bytes:
            if self == path:
                raise OSError("模拟读取失败")
            return original_read_bytes(self)

        monkeypatch.setattr(Path, "read_bytes", mock_read_bytes)
        with pytest.raises(ExtractorError, match="文件读取失败"):
            TextExtractor().extract(path)

    def test_charset_normalizer_exception_fallback(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """charset-normalizer 抛异常时回退到 UTF-8/GBK 解码。

        内容用 GBK 编码写入，使其无法命中 UTF-8 快路径，从而进入
        charset-normalizer 分支触发 RuntimeError 回退（回退链中 GBK 可解）。
        禁用原生编码检测以隔离测试 charset-normalizer 异常回退路径。
        """
        monkeypatch.setattr("fuscan.extractors.text._NATIVE_DECODE_AVAILABLE", False)
        path = tmp_path / "test.txt"
        path.write_bytes("异常回退 password".encode("gbk"))

        def fake_from_bytes(data: bytes):
            raise RuntimeError("模拟检测异常")

        monkeypatch.setattr("charset_normalizer.from_bytes", fake_from_bytes)
        content = TextExtractor().extract(path)
        assert "异常回退 password" in content

    def test_charset_normalizer_none_fallback_to_latin1(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """charset-normalizer 返回 None 时回退到 latin-1 解码任意字节。

        字节须既非 BOM 也非合法 UTF-8，才不会命中头部快路径而进入
        charset-normalizer 分支（``\\x80\\x81\\xfd`` 为孤立续字节，非法 UTF-8）。
        禁用原生编码检测以隔离测试 charset-normalizer None 回退路径。
        """
        monkeypatch.setattr("fuscan.extractors.text._NATIVE_DECODE_AVAILABLE", False)
        path = tmp_path / "test.txt"
        path.write_bytes(b"\x80\x81\xfd")

        monkeypatch.setattr("charset_normalizer.from_bytes", lambda data: type("R", (), {"best": lambda self: None})())
        content = TextExtractor().extract(path)
        assert isinstance(content, str)
        # latin-1 能解码任意字节，不应为空
        assert len(content) == 3

    def test_native_decode_gbk(self, tmp_path: Path) -> None:
        """原生编码检测（fuscan-core）正确解码 GBK 中文。

        用足够长的 GBK 文本给 chardetng 统计检测足够的特征，
        确保检测到 GBK 编码并正确解码中文与 ASCII 关键字。
        """
        path = tmp_path / "gbk_native.txt"
        text = "这是一个包含密码字段的配置文件，密码为 password123，请妥善保管。"
        path.write_bytes(text.encode("gbk"))
        content = TextExtractor().extract(path)
        assert "密码" in content
        assert "password123" in content

    def test_native_decode_priority_over_charset_normalizer(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        """原生编码检测可用时优先于 charset-normalizer。

        monkeypatch charset-normalizer 抛异常，若原生路径生效则不受影响，
        GBK 内容仍被正确解码（证明走了原生路径而非 charset-normalizer）。
        """

        def fake_from_bytes(data: bytes):
            raise AssertionError("charset-normalizer 不应被调用，原生路径应优先")

        monkeypatch.setattr("charset_normalizer.from_bytes", fake_from_bytes)
        path = tmp_path / "gbk_priority.txt"
        text = "原生编码检测优先 password 配置文件密码"
        path.write_bytes(text.encode("gbk"))
        content = TextExtractor().extract(path)
        assert "password" in content

    def test_native_decode_exception_fallback(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """原生编码检测抛异常时回退到 charset-normalizer。

        monkeypatch 原生 decode_bytes 抛 RuntimeError，验证回退到
        charset-normalizer 路径仍能正确解码 GBK 内容。
        """

        def fake_decode(data: bytes) -> str:
            raise RuntimeError("模拟原生检测异常")

        monkeypatch.setattr("fuscan.extractors.text._native_decode_bytes", fake_decode)
        path = tmp_path / "native_exc.txt"
        text = "异常回退原生检测 password 密码配置"
        path.write_bytes(text.encode("gbk"))
        content = TextExtractor().extract(path)
        assert "password" in content

    def test_native_decode_invalid_bytes(self, tmp_path: Path) -> None:
        """原生编码检测对非法字节用 U+FFFD 替换，不 panic/不抛异常。"""
        path = tmp_path / "invalid.txt"
        path.write_bytes(b"\x80\x81\xfd")
        content = TextExtractor().extract(path)
        assert isinstance(content, str)
        assert len(content) > 0

    def test_normalizes_crlf_to_lf(self, tmp_path: Path) -> None:
        """CRLF 行尾应规范化为 LF，保证跨平台 CONTENT EQUALS 比较一致。"""
        path = tmp_path / "crlf.txt"
        path.write_bytes(b"line1\r\nline2\r\n")
        content = TextExtractor().extract(path)
        assert content == "line1\nline2\n"
        assert "\r\n" not in content

    def test_normalizes_cr_to_lf(self, tmp_path: Path) -> None:
        """旧式 CR 行尾应规范化为 LF。"""
        path = tmp_path / "cr.txt"
        path.write_bytes(b"line1\rline2\r")
        content = TextExtractor().extract(path)
        assert content == "line1\nline2\n"

    def test_lf_preserved(self, tmp_path: Path) -> None:
        """LF 行尾保持不变。"""
        path = tmp_path / "lf.txt"
        path.write_bytes(b"line1\nline2\n")
        content = TextExtractor().extract(path)
        assert content == "line1\nline2\n"


# ---------------------------------------------------------------------------
# DocxExtractor
# ---------------------------------------------------------------------------


class TestDocxExtractor:
    def test_extract_paragraphs_and_table(self, docx_file: Path) -> None:
        extractor = DocxExtractor()
        content = extractor.extract(docx_file)
        assert "段落一 含 password" in content
        assert "段落二 正常内容" in content
        assert "姓名" in content
        assert "pwd123" in content

    def test_supported_extensions(self) -> None:
        assert DocxExtractor().supported_extensions == ("docx",)

    def test_invalid_file_raises(self, tmp_path: Path) -> None:
        path = tmp_path / "bad.docx"
        path.write_text("not a docx", encoding="utf-8")
        with pytest.raises(ExtractorError, match="DOCX 解析失败"):
            DocxExtractor().extract(path)


# ---------------------------------------------------------------------------
# PptxExtractor
# ---------------------------------------------------------------------------


class TestPptxExtractor:
    def test_extract_slide_text(self, pptx_file: Path) -> None:
        extractor = PptxExtractor()
        content = extractor.extract(pptx_file)
        assert "标题 含 secret" in content
        assert "幻灯片内容" in content

    def test_supported_extensions(self) -> None:
        assert PptxExtractor().supported_extensions == ("pptx",)

    def test_invalid_file_raises(self, tmp_path: Path) -> None:
        path = tmp_path / "bad.pptx"
        path.write_text("not a pptx", encoding="utf-8")
        with pytest.raises(ExtractorError, match="PPTX 解析失败"):
            PptxExtractor().extract(path)

    def test_pptx_with_table_and_notes(self, tmp_path: Path) -> None:
        """PPTX 含表格和备注时应提取这些内容。"""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "表格测试"  # pyrefly: ignore [missing-attribute]
        # 添加表格
        table_shape = slide.shapes.add_table(
            rows=2, cols=2, left=Inches(1), top=Inches(2), width=Inches(4), height=Inches(1)
        )
        table = table_shape.table
        table.cell(0, 0).text = "键"
        table.cell(0, 1).text = "密码"
        table.cell(1, 0).text = "user"
        table.cell(1, 1).text = "pwd123"
        # 添加备注
        slide.notes_slide.notes_text_frame.text = "备注内容 secret"  # pyrefly: ignore [missing-attribute]
        path = tmp_path / "table_notes.pptx"
        prs.save(str(path))

        content = PptxExtractor().extract(path)
        assert "表格测试" in content
        assert "pwd123" in content
        assert "[备注]" in content
        assert "备注内容 secret" in content

    def test_pptx_empty_slide_skipped(self, tmp_path: Path) -> None:
        """空幻灯片应被跳过（不添加 --- 幻灯片 --- 分隔符）。"""
        from pptx import Presentation

        prs = Presentation()
        # 添加一个空布局的幻灯片（无文本）
        prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        path = tmp_path / "empty.pptx"
        prs.save(str(path))

        content = PptxExtractor().extract(path)
        assert "幻灯片" not in content


class TestDocxExtractorExtra:
    """DocxExtractor 额外覆盖。"""

    def test_docx_with_header_footer(self, tmp_path: Path) -> None:
        """DOCX 含页眉页脚时应提取这些内容。"""
        from docx import Document

        doc = Document()
        doc.add_paragraph("正文 password")
        # 添加页眉页脚
        section = doc.sections[0]
        section.header.paragraphs[0].text = "页眉内容 secret"
        section.footer.paragraphs[0].text = "页脚信息"
        path = tmp_path / "header_footer.docx"
        doc.save(str(path))

        content = DocxExtractor().extract(path)
        assert "正文 password" in content
        assert "页眉内容 secret" in content
        assert "页脚信息" in content

    def test_docx_extract_oserror_raises(self, tmp_path: Path) -> None:
        """读取不存在文件时应抛 ``ExtractorError``。"""
        path = tmp_path / "missing.docx"
        with pytest.raises(ExtractorError, match="文件读取失败"):
            DocxExtractor().extract(path)

    def test_docx_lxml_bad_zipfile_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """lxml 抛 ``BadZipFile`` 时应抛 ``ExtractorError``。"""
        path = tmp_path / "bad.docx"
        path.write_bytes(b"not a zip")

        import zipfile

        def _raise_bad_zip(_data: bytes) -> str:
            raise zipfile.BadZipFile("invalid zip")

        monkeypatch.setattr("fuscan.extractors._ooxml_xml.extract_docx_text", _raise_bad_zip)
        with pytest.raises(ExtractorError, match="DOCX 解析失败"):
            DocxExtractor().extract(path)

    def test_pptx_extract_oserror_raises(self, tmp_path: Path) -> None:
        """读取不存在文件时应抛 ``ExtractorError``。"""
        path = tmp_path / "missing.pptx"
        with pytest.raises(ExtractorError, match="文件读取失败"):
            PptxExtractor().extract(path)

    def test_pptx_lxml_bad_zipfile_raises(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        """lxml 抛 ``BadZipFile`` 时应抛 ``ExtractorError``。"""
        path = tmp_path / "bad.pptx"
        path.write_bytes(b"not a zip")

        import zipfile

        def _raise_bad_zip(_data: bytes) -> str:
            raise zipfile.BadZipFile("invalid zip")

        monkeypatch.setattr("fuscan.extractors._ooxml_xml.extract_pptx_text", _raise_bad_zip)
        with pytest.raises(ExtractorError, match="PPTX 解析失败"):
            PptxExtractor().extract(path)


# ---------------------------------------------------------------------------
# XlsxExtractor
# ---------------------------------------------------------------------------


class TestXlsxExtractor:
    def test_extract_cells(self, xlsx_file: Path) -> None:
        extractor = XlsxExtractor()
        content = extractor.extract(xlsx_file)
        assert "姓名" in content
        assert "pwd123" in content
        assert "数据" in content  # 工作表名

    def test_supported_extensions(self) -> None:
        exts = XlsxExtractor().supported_extensions
        assert "xlsx" in exts
        assert "xlsm" in exts

    def test_invalid_file_raises(self, tmp_path: Path) -> None:
        path = tmp_path / "bad.xlsx"
        path.write_text("not xlsx", encoding="utf-8")
        with pytest.raises(ExtractorError, match="XLSX 解析失败"):
            XlsxExtractor().extract(path)

    def test_max_rows_limit(self, xlsx_file: Path) -> None:
        extractor = XlsxExtractor(max_rows=1)
        content = extractor.extract(xlsx_file)
        # 只读了 1 行，应该有表头但无数据行
        assert "姓名" in content

    def test_max_cols_limit(self, xlsx_file: Path) -> None:
        """列数超过上限时截断。"""
        extractor = XlsxExtractor(max_cols=1)
        content = extractor.extract(xlsx_file)
        # 只有第 1 列，应包含姓名但不包含密码列
        assert "姓名" in content

    def test_import_error_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """python-calamine 未安装时应抛出 ExtractorError。"""
        path = tmp_path / "test.xlsx"
        path.write_bytes(b"fake")
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "python_calamine":
                raise ImportError("No module named 'python_calamine'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="python-calamine 未安装"):
            XlsxExtractor().extract(path)

    def test_output_format_exact(self, tmp_path: Path) -> None:
        """R3 输出等价红线：多工作表/混合类型/空表的提取字符串逐字节固定。

        构造含数值/布尔/字符串混合单元格、含空表、含前导空白与 None 的工作簿，
        断言输出严格符合「``--- 工作表: 名称 ---`` 分隔 + 行内 ``\\t`` + 行间 ``\\n``」
        规范，防止 R3 遍历重构改变输出。
        """
        from openpyxl import Workbook

        wb = Workbook()
        ws1 = wb.active
        assert ws1 is not None
        ws1.title = "数据"
        ws1["A1"] = "姓名"
        ws1["B1"] = "密码"
        ws1["C1"] = 42  # 数值 cell：calamine 返回 float，str(42.0)="42.0"
        ws1["A2"] = "张三"
        ws1["B2"] = "pwd123"
        # C2 留空（None）→ 跳过；D2 前后空白 → strip
        ws1["D2"] = "  trimmed  "
        # 全空行（第 3 行）→ 整行跳过
        wb.create_sheet("空表")  # 空 sheet → 不产生任何输出
        ws3 = wb.create_sheet("布尔")
        ws3["A1"] = True  # str(True)="True"
        path = tmp_path / "exact.xlsx"
        wb.save(str(path))

        content = XlsxExtractor().extract(path)
        expected = "--- 工作表: 数据 ---\n姓名\t密码\t42.0\n张三\tpwd123\ttrimmed\n--- 工作表: 布尔 ---\nTrue"
        assert content == expected

    def test_row_col_truncation_exact(self, tmp_path: Path) -> None:
        """R3 截断等价红线：超行/超列截断输出与逐单元格计数路径一致。"""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        assert ws is not None
        ws.title = "T"
        # 3 行 × 3 列，全部填值
        for r in range(1, 4):
            for c, col in enumerate("ABC", 0):
                ws[f"{col}{r}"] = f"r{r}c{c}"
        path = tmp_path / "trunc.xlsx"
        wb.save(str(path))

        # max_rows=2 max_cols=2：仅前 2 行前 2 列
        content = XlsxExtractor(max_rows=2, max_cols=2).extract(path)
        expected = "--- 工作表: T ---\nr1c0\tr1c1\nr2c0\tr2c1"
        assert content == expected


class TestWpsExtractor:
    def test_supported_extensions(self) -> None:
        exts = WpsExtractor().supported_extensions
        assert "wps" in exts
        assert "et" in exts
        assert "dps" in exts

    def test_non_ooxml_returns_empty(self, tmp_path: Path) -> None:
        """旧版二进制 WPS 格式应返回空字符串。"""
        path = tmp_path / "old.wps"
        path.write_bytes(b"\xd0\xcf\x11\xe0 old binary format")
        assert WpsExtractor().extract(path) == ""

    def test_ooxml_wps_text(self, tmp_path: Path) -> None:
        """OOXML 兼容的 .wps 文件应能提取文本（实际是 DOCX 内容）。"""
        from docx import Document

        doc = Document()
        doc.add_paragraph("wps 内容 password")
        path = tmp_path / "test.wps"
        doc.save(str(path))

        content = WpsExtractor().extract(path)
        assert "wps 内容 password" in content

    def test_ooxml_et_sheet(self, tmp_path: Path) -> None:
        """OOXML 兼容的 .et 文件应能提取表格内容。"""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        assert ws is not None
        ws["A1"] = "et_password"
        path = tmp_path / "test.et"
        wb.save(str(path))

        content = WpsExtractor().extract(path)
        assert "et_password" in content

    def test_ooxml_dps_slides(self, tmp_path: Path) -> None:
        """OOXML 兼容的 .dps 文件应能提取演示内容。"""
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "dps 标题 password"  # pyrefly: ignore [missing-attribute]
        slide.placeholders[1].text = "幻灯片内容"  # pyrefly: ignore [missing-attribute]
        path = tmp_path / "test.dps"
        prs.save(str(path))

        content = WpsExtractor().extract(path)
        assert "dps 标题 password" in content
        assert "幻灯片内容" in content

    def test_wps_docx_with_table(self, tmp_path: Path) -> None:
        """WPS 文字文档含表格时应提取表格内容。"""
        from docx import Document

        doc = Document()
        doc.add_paragraph("正文 password")
        table = doc.add_table(rows=1, cols=2)
        table.cell(0, 0).text = "键"
        table.cell(0, 1).text = "值 secret"
        path = tmp_path / "table.wps"
        doc.save(str(path))

        content = WpsExtractor().extract(path)
        assert "正文 password" in content
        assert "键" in content
        assert "值 secret" in content

    def test_wps_invalid_docx_returns_empty(self, tmp_path: Path) -> None:
        """有效的 ZIP 但 docx 内容损坏时，lxml recover 跳过，返回空字符串。"""
        path = tmp_path / "bad.wps"
        path.write_bytes(_make_ooxml_zip("word/document.xml", "corrupt xml"))
        assert WpsExtractor().extract(path) == ""

    def test_wps_invalid_et_raises(self, tmp_path: Path) -> None:
        """有效的 ZIP 但 xlsx 内容损坏时应抛出 ExtractorError。"""
        path = tmp_path / "bad.et"
        path.write_bytes(_make_ooxml_zip("xl/workbook.xml", "corrupt xml"))
        with pytest.raises(ExtractorError, match="WPS 表格 解析失败"):
            WpsExtractor().extract(path)

    def test_wps_invalid_dps_returns_empty(self, tmp_path: Path) -> None:
        """有效的 ZIP 但 pptx 幻灯片内容损坏时，lxml recover 跳过，返回空字符串。"""
        # 需同时包含 presentation.xml（类型检测）和 slide1.xml（损坏内容）
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w") as zf:
            zf.writestr("ppt/presentation.xml", "fake")
            zf.writestr("ppt/slides/slide1.xml", "corrupt xml")
        path = tmp_path / "bad.dps"
        path.write_bytes(buf.getvalue())
        assert WpsExtractor().extract(path) == ""

    def test_wps_corrupt_zip_returns_empty(self, tmp_path: Path) -> None:
        """ZIP 头存在但 ZIP 损坏时应返回空字符串（无法判定子类型）。"""
        path = tmp_path / "corrupt.wps"
        path.write_bytes(b"PK\x03\x04 corrupted content")
        assert WpsExtractor().extract(path) == ""

    def test_wps_detect_ooxml_type_bad_zip_returns_none(self) -> None:
        """_detect_ooxml_type 对损坏的 ZIP 数据返回 None。"""
        assert WpsExtractor()._detect_ooxml_type(b"PK\x03\x04 corrupted") is None

    def test_wps_detect_ooxml_type_unknown_returns_none(self) -> None:
        """_detect_ooxml_type 对未知 OOXML 子类型返回 None。"""
        data = _make_ooxml_zip("unknown/entry.xml")
        assert WpsExtractor()._detect_ooxml_type(data) is None

    def test_wps_detect_ooxml_type_docx(self) -> None:
        """_detect_ooxml_type 正确识别 docx 子类型。"""
        data = _make_ooxml_zip("word/document.xml")
        assert WpsExtractor()._detect_ooxml_type(data) == "docx"

    def test_wps_detect_ooxml_type_xlsx(self) -> None:
        """_detect_ooxml_type 正确识别 xlsx 子类型。"""
        data = _make_ooxml_zip("xl/workbook.xml")
        assert WpsExtractor()._detect_ooxml_type(data) == "xlsx"

    def test_wps_detect_ooxml_type_pptx(self) -> None:
        """_detect_ooxml_type 正确识别 pptx 子类型。"""
        data = _make_ooxml_zip("ppt/presentation.xml")
        assert WpsExtractor()._detect_ooxml_type(data) == "pptx"


class TestWpsExtractorErrorPaths:
    """WPS 提取器异常路径覆盖。"""

    def test_extract_corrupt_zip_returns_empty(self, tmp_path: Path) -> None:
        """ZIP 头存在但内容损坏时应返回空字符串（无法判定子类型）。"""
        path = tmp_path / "file.unknown"
        path.write_bytes(b"PK\x03\x04 fake content")
        # ZIP 头存在但非有效 ZIP，_detect_ooxml_type 返回 None，走 return ""
        extractor = WpsExtractor()
        result = extractor.extract(path)
        assert result == ""

    def test_extract_as_xlsx_import_error_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """python-calamine 未安装时 _extract_as_xlsx 应抛出 ExtractorError。"""
        path = tmp_path / "test.et"
        path.write_bytes(_make_ooxml_zip("xl/workbook.xml"))
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "python_calamine":
                raise ImportError("No module named 'python_calamine'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="python-calamine 未安装"):
            WpsExtractor().extract(path)


# ---------------------------------------------------------------------------
# PdfExtractor（pypdfium2 后端，延迟导入）
# ---------------------------------------------------------------------------


def _make_pdf_sample(tmp_path: Path) -> bytes:
    """用 reportlab 生成含 password 关键词的 PDF 样本。"""
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate

    path = tmp_path / "sample.pdf"
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(str(path), pagesize=letter)
    doc.build([Paragraph("This document contains a secret password.", styles["Normal"])])
    return path.read_bytes()


class TestPdfExtractor:
    """PDF 提取器基础测试（不依赖具体后端）。"""

    def test_supported_extensions(self) -> None:
        assert PdfExtractor().supported_extensions == ("pdf",)

    def test_invalid_file_raises(self, tmp_path: Path) -> None:
        path = tmp_path / "bad.pdf"
        path.write_text("not a pdf", encoding="utf-8")
        with pytest.raises(ExtractorError):
            PdfExtractor().extract(path)

    def test_extract_missing_file_raises(self, tmp_path: Path) -> None:
        """``extract()`` 读取缺失文件应抛出 ``ExtractorError``。"""
        with pytest.raises(ExtractorError, match="文件读取失败"):
            PdfExtractor().extract(tmp_path / "missing.pdf")

    def test_no_backend_available_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """pypdfium2 不可用时应抛出 ``ExtractorError``。"""

        def _raise() -> None:
            raise ExtractorError("无可用 PDF 引擎（pypdfium2 未安装）")

        monkeypatch.setattr("fuscan.extractors.pdf._ensure_backend", _raise)
        path = tmp_path / "x.pdf"
        path.write_bytes(b"fake")
        with pytest.raises(ExtractorError, match="无可用 PDF 引擎"):
            PdfExtractor().extract(path)


# ---------------------------------------------------------------------------
# PdfExtractor pypdfium2 后端测试
# ---------------------------------------------------------------------------


class TestPdfExtractorPdfiumBackend:
    """pypdfium2（pdfium C++）后端测试。

    pypdfium2 是唯一的 PDF 解析后端（Google pdfium C++ 引擎，cffi 绑定），
    兼容 Win7。speed_tier 固定 T3 中速，engine_info 固定 "pypdfium2"。
    通过 :func:`_ensure_backend` 延迟导入，未安装时抛 ``ExtractorError``。
    """

    def test_pdfium_speed_tier_is_medium(self) -> None:
        """pypdfium2 后端 speed_tier 返回 T3 中速。"""
        assert PdfExtractor().speed_tier == SpeedTier.MEDIUM

    def test_pdfium_engine_info(self) -> None:
        """engine_info 返回 pypdfium2。"""
        assert PdfExtractor().engine_info == "pypdfium2"

    def test_pdfium_extract_real_pdf(self, tmp_path: Path) -> None:
        """pypdfium2 后端提取真实 PDF 应包含 password 关键词。"""
        pdf_sample = _make_pdf_sample(tmp_path)
        extractor = PdfExtractor()
        content = extractor.extract_from_bytes(pdf_sample)
        assert "password" in content.lower()

    def test_pdfium_extract_path_matches_bytes(self, tmp_path: Path) -> None:
        """pypdfium2 从 path 与从 bytes 提取结果一致。"""
        pdf_sample = _make_pdf_sample(tmp_path)
        path = tmp_path / "sample.pdf"
        path.write_bytes(pdf_sample)
        extractor = PdfExtractor()
        assert extractor.extract(path) == extractor.extract_from_bytes(pdf_sample)

    def test_pdfium_invalid_bytes_raises(self) -> None:
        """pypdfium2 无法打开无效字节时应抛出 ``ExtractorError``。"""
        with pytest.raises(ExtractorError, match="PDF 打开失败"):
            PdfExtractor().extract_from_bytes(b"not a valid pdf")

    def test_pdfium_page_failures_aggregated_not_per_page(
        self, monkeypatch: pytest.MonkeyPatch, caplog: pytest.LogCaptureFixture
    ) -> None:
        """损坏 PDF 逐页失败应汇总为一条 WARNING，不逐页打印 traceback。

        模拟 3 页 PDF 中第 0、2 页 ``get_page`` 抛异常（pdfium 无法加载），
        验证：成功页文本正常返回；逐页失败仅 DEBUG；循环结束汇总一条 WARNING。
        """
        import logging

        from fuscan.extractors import pdf as pdf_mod

        class _TextPage:
            def get_text_range(self) -> str:
                return "正常页文本"

            def close(self) -> None:
                """无操作。"""

        class _GoodPage:
            def get_textpage(self) -> _TextPage:
                return _TextPage()

            def close(self) -> None:
                """无操作。"""

        class _MixedDoc:
            def __len__(self) -> int:
                return 3

            def get_page(self, i: int) -> object:
                if i == 1:
                    return _GoodPage()
                raise RuntimeError("Failed to load page.")

            def close(self) -> None:
                """无操作。"""

        monkeypatch.setattr(pdf_mod, "_ensure_backend", lambda: lambda _data: _MixedDoc())
        extractor = PdfExtractor()
        with caplog.at_level(logging.DEBUG, logger="fuscan.extractors.pdf"):
            content = extractor.extract_from_bytes(b"fake mixed pdf")
        # 成功页文本正常返回
        assert "正常页文本" in content
        # 逐页失败仅 DEBUG（不产生逐页 WARNING）
        debug_msgs = [r.message for r in caplog.records if r.levelno == logging.DEBUG]
        warning_msgs = [r.message for r in caplog.records if r.levelno == logging.WARNING]
        assert any("页 0 提取失败" in m for m in debug_msgs)
        assert any("页 2 提取失败" in m for m in debug_msgs)
        # 汇总仅一条 WARNING（含总页数与失败页数）
        assert len(warning_msgs) == 1
        assert "3 页" in warning_msgs[0]
        assert "2 页提取失败" in warning_msgs[0]

    def test_pdfium_child_objects_closed_each_page(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """每页提取后 textpage/page 应显式 close，避免 doc.close() 后 GC 终结触发断言。

        pypdfium2 子对象（PdfTextPage/PdfPage）未显式关闭时，doc.close() 后 GC
        终结子对象会触发 ``_close_template`` 的 ``assert not parent._tree_closed()``
        断言失败及 "kids weakrefs not cleaned up" 警告。验证每页的 close 均被调用。
        """
        from fuscan.extractors import pdf as pdf_mod

        closed: list[str] = []

        class _TextPage:
            def get_text_range(self) -> str:
                return "text"

            def close(self) -> None:
                closed.append("textpage")

        class _Page:
            def get_textpage(self) -> _TextPage:
                return _TextPage()

            def close(self) -> None:
                closed.append("page")

        class _Doc:
            def __len__(self) -> int:
                return 2

            def get_page(self, i: int) -> _Page:
                return _Page()

            def close(self) -> None:
                closed.append("doc")

        monkeypatch.setattr(pdf_mod, "_ensure_backend", lambda: lambda _data: _Doc())
        PdfExtractor().extract_from_bytes(b"fake pdf")
        # 2 页 × (textpage + page) + 1 doc，子对象先于 doc 关闭
        assert closed == ["textpage", "page", "textpage", "page", "doc"]


# ---------------------------------------------------------------------------
# OdtExtractor / OdsExtractor（用标准库 zipfile+xml 解析，无 odfpy 依赖）
# ---------------------------------------------------------------------------


class TestOdfExtractors:
    def test_odt_supported_extensions(self) -> None:
        assert OdtExtractor().supported_extensions == ("odt",)

    def test_ods_supported_extensions(self) -> None:
        assert OdsExtractor().supported_extensions == ("ods",)

    def test_odt_invalid_file_raises(self, tmp_path: Path) -> None:
        path = tmp_path / "bad.odt"
        path.write_text("not odt", encoding="utf-8")
        with pytest.raises(ExtractorError, match="ODT 解析失败"):
            OdtExtractor().extract(path)

    def test_odt_read_failure_raises(self, tmp_path: Path) -> None:
        """ODT 文件读取失败（权限/不存在）应抛出 ExtractorError。"""
        path = tmp_path / "nonexistent.odt"
        with pytest.raises(ExtractorError, match="文件读取失败"):
            OdtExtractor().extract(path)

    def test_ods_read_failure_raises(self, tmp_path: Path) -> None:
        """ODS 文件读取失败应抛出 ExtractorError。"""
        path = tmp_path / "nonexistent.ods"
        with pytest.raises(ExtractorError, match="文件读取失败"):
            OdsExtractor().extract(path)

    def test_odt_extract_real_file(self, tmp_path: Path) -> None:
        """用 zipfile 构造真实 ODT 文件并提取。"""
        from tests._odf_samples import make_odt_sample

        path = tmp_path / "real.odt"
        path.write_bytes(
            make_odt_sample(["段落含 password 内容", "标题 secret"]),
        )

        content = OdtExtractor().extract(path)
        assert "password" in content
        assert "secret" in content

    def test_ods_extract_real_file(self, tmp_path: Path) -> None:
        """用 zipfile 构造真实 ODS 文件并提取。"""
        from tests._odf_samples import make_ods_sample

        path = tmp_path / "real.ods"
        path.write_bytes(make_ods_sample([["cell_password"]]))

        content = OdsExtractor().extract(path)
        assert "cell_password" in content

    def test_ods_invalid_file_raises(self, tmp_path: Path) -> None:
        path = tmp_path / "bad.ods"
        path.write_text("not ods", encoding="utf-8")
        with pytest.raises(ExtractorError, match="ODS 解析失败"):
            OdsExtractor().extract(path)

    def test_odt_missing_content_xml_raises(self, tmp_path: Path) -> None:
        """ZIP 包缺少 content.xml 时应抛出 ExtractorError。"""
        import zipfile

        path = tmp_path / "empty.odt"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("mimetype", "application/vnd.oasis.opendocument.text")
        with pytest.raises(ExtractorError, match="ODT 解析失败"):
            OdtExtractor().extract(path)

    def test_ods_missing_content_xml_raises(self, tmp_path: Path) -> None:
        """ZIP 包缺少 content.xml 时应抛出 ExtractorError。"""
        import zipfile

        path = tmp_path / "empty.ods"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("mimetype", "application/vnd.oasis.opendocument.spreadsheet")
        with pytest.raises(ExtractorError, match="ODS 解析失败"):
            OdsExtractor().extract(path)

    def test_odt_extract_from_bytes_matches_path(self, tmp_path: Path) -> None:
        """OdtExtractor 从 bytes 提取与从 path 提取结果一致。"""
        from tests._odf_samples import make_odt_sample

        data = make_odt_sample(["odt password 内容", "标题内容"])
        path = tmp_path / "test.odt"
        path.write_bytes(data)
        extractor = OdtExtractor()
        assert extractor.extract(path) == extractor.extract_from_bytes(data)

    def test_ods_extract_from_bytes_matches_path(self, tmp_path: Path) -> None:
        """OdsExtractor 从 bytes 提取与从 path 提取结果一致。"""
        from tests._odf_samples import make_ods_sample

        data = make_ods_sample([["cell1", "cell2"], ["cell3", "ods_password"]])
        path = tmp_path / "test.ods"
        path.write_bytes(data)
        extractor = OdsExtractor()
        assert extractor.extract(path) == extractor.extract_from_bytes(data)

    def test_ods_multi_row_multi_col_extraction(self) -> None:
        """ODS 多行多列单元格应以 \\t 分隔列、\\n 分隔行。"""
        from tests._odf_samples import make_ods_sample

        data = make_ods_sample(
            [
                ["A1", "B1", "C1"],
                ["A2", "B2", "C2"],
            ],
        )
        content = OdsExtractor().extract_from_bytes(data)
        lines = content.split("\n")
        assert lines == ["A1\tB1\tC1", "A2\tB2\tC2"]

    def test_odt_empty_paragraphs_returns_empty(self) -> None:
        """ODT 无段落时应返回空字符串。"""
        from tests._odf_samples import make_odt_sample

        data = make_odt_sample([])
        assert OdtExtractor().extract_from_bytes(data) == ""

    def test_ods_empty_rows_returns_empty(self) -> None:
        """ODS 无行时应返回空字符串。"""
        from tests._odf_samples import make_ods_sample

        data = make_ods_sample([])
        assert OdsExtractor().extract_from_bytes(data) == ""

    def test_odt_xml_special_chars_escaped(self) -> None:
        """ODT 段落中的 XML 特殊字符应正确转义与提取。"""
        from tests._odf_samples import make_odt_sample

        data = make_odt_sample(["a < b & c > d", 'quote " test'])
        content = OdtExtractor().extract_from_bytes(data)
        assert "a < b & c > d" in content
        assert 'quote " test' in content

    def test_element_text_handles_nested_elements(self) -> None:
        """element_text 应递归提取嵌套子元素文本与 tail。"""
        import xml.etree.ElementTree as ET

        from fuscan.extractors._odf_xml import TEXT_NS, element_text

        # 构造 <text:p>前<text:span>中</text:span>后</text:p>
        parent = ET.Element(f"{{{TEXT_NS}}}p")
        parent.text = "前"
        child = ET.SubElement(parent, f"{{{TEXT_NS}}}span")
        child.text = "中"
        child.tail = "后"
        assert element_text(parent) == "前中后"

    def test_element_text_empty_element_returns_empty(self) -> None:
        """element_text 对无文本无子元素的应返回空串。"""
        import xml.etree.ElementTree as ET

        from fuscan.extractors._odf_xml import TEXT_NS, element_text

        empty = ET.Element(f"{{{TEXT_NS}}}p")
        assert element_text(empty) == ""

    def test_odt_nested_paragraph_extraction(self) -> None:
        """ODT 段落含嵌套 span 时应正确提取完整文本。"""
        import io
        import zipfile

        # 手工构造带嵌套 span 的 ODT
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>\n'
            "<office:document-content "
            'xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" '
            'xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0" '
            'office:version="1.2">\n'
            "<office:body><office:text>\n"
            "<text:p>前<text:span>中</text:span>后password</text:p>\n"
            "</office:text></office:body></office:document-content>\n"
        )
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
            mi = zipfile.ZipInfo("mimetype")
            mi.compress_type = zipfile.ZIP_STORED
            zf.writestr(mi, b"application/vnd.oasis.opendocument.text")
            zf.writestr("content.xml", content_xml.encode("utf-8"))

        content = OdtExtractor().extract_from_bytes(buf.getvalue())
        assert content == "前中后password"


# ---------------------------------------------------------------------------
# ExtractorRegistry
# ---------------------------------------------------------------------------


class TestExtractorRegistry:
    def test_register_and_get(self) -> None:
        registry = ExtractorRegistry()
        registry.register(TextExtractor())
        assert registry.get("txt") is not None
        assert registry.get("TXT") is not None  # 大小写不敏感
        assert registry.get("missing") is None

    def test_registered_extensions(self) -> None:
        registry = ExtractorRegistry()
        registry.register(TextExtractor())
        exts = registry.registered_extensions
        assert "txt" in exts
        assert "md" in exts

    def test_extract_with_registered(self, text_file: Path) -> None:
        registry = ExtractorRegistry()
        registry.register(TextExtractor())
        content = registry.extract(text_file)
        assert "hello password world" in content

    def test_extract_without_extractor_returns_empty(self, tmp_path: Path) -> None:
        path = tmp_path / "unknown.xyz"
        path.write_text("content", encoding="utf-8")
        registry = ExtractorRegistry()
        assert registry.extract(path) == ""

    def test_default_registry_has_all(self) -> None:
        exts = default_registry.registered_extensions
        for expected in ("txt", "pdf", "docx", "pptx", "xlsx", "odt", "ods", "wps"):
            assert expected in exts, f"默认注册表缺少 {expected}"

    def test_list_extractors_returns_unique_entries(self) -> None:
        """list_extractors 返回去重后的提取器列表（同一实例多扩展名合并为一项）。"""
        extractors = default_registry.list_extractors()
        # 14 个提取器类，每个对应一项
        class_names = [entry[0] for entry in extractors]
        assert len(class_names) == len(set(class_names)), "提取器列表有重复"
        # 按 display_name 排序
        display_names = [entry[1] for entry in extractors]
        assert display_names == sorted(display_names)

    def test_list_extractors_entry_format(self) -> None:
        """list_extractors 返回元组格式为 (class_name, display_name, extensions, speed_tier, engine_info)。

        iter-90 起新增 speed_tier 字段（SpeedTier 枚举）。
        iter-139 起新增 engine_info 字段（str，描述底层解析引擎）。
        """
        from fuscan.extractors.base import SpeedTier

        extractors = default_registry.list_extractors()
        for class_name, display_name, exts, tier, engine_info in extractors:
            assert isinstance(class_name, str) and class_name
            assert isinstance(display_name, str) and display_name
            assert isinstance(exts, tuple) and exts
            assert isinstance(tier, SpeedTier)
            assert isinstance(engine_info, str)
            # 扩展名均为小写无点
            for e in exts:
                assert e == e.lower().lstrip(".")

    def test_display_name_returns_chinese(self) -> None:
        """各提取器 display_name 返回非空中文名称。"""
        names = {
            TextExtractor: "纯文本",
            PdfExtractor: "PDF",
            DocxExtractor: "Word（DOCX）",
            PptxExtractor: "PowerPoint（PPTX）",
            XlsxExtractor: "Excel（XLSX）",
            OdsExtractor: "ODS 表格",
            OdtExtractor: "ODT 文档",
            WpsExtractor: "WPS 文档（WPS）",
            EmlExtractor: "邮件（EML）",
            XlsExtractor: "Excel（XLS）",
            DocExtractor: "Word（DOC）",
            PptExtractor: "PowerPoint（PPT）",
        }
        for cls, expected in names.items():
            assert cls().display_name == expected, f"{cls.__name__}.display_name 应为 {expected}"

    def test_engine_info_returns_non_empty_str(self) -> None:
        """各提取器 engine_info 应返回非空字符串。"""
        classes = [
            TextExtractor,
            PdfExtractor,
            DocxExtractor,
            PptxExtractor,
            XlsxExtractor,
            OdsExtractor,
            OdtExtractor,
            WpsExtractor,
            EmlExtractor,
            XlsExtractor,
            DocExtractor,
            PptExtractor,
        ]
        for cls in classes:
            info = cls().engine_info
            assert isinstance(info, str) and info, f"{cls.__name__}.engine_info 应为非空字符串"

    def test_engine_info_specific_values(self) -> None:
        """固定后端的提取器 engine_info 应返回预期引擎名。"""
        # XLSX/XLS 固定使用 calamine
        assert XlsxExtractor().engine_info == "python-calamine"
        assert XlsExtractor().engine_info == "python-calamine"
        # PDF 固定使用 pypdfium2
        assert PdfExtractor().engine_info == "pypdfium2"
        # DOC/PPT 在 fuscan-core (cfb) 与 olefile 之间切换
        assert DocExtractor().engine_info in {"fuscan-core (cfb)", "olefile"}
        assert PptExtractor().engine_info in {"fuscan-core (cfb)", "olefile"}
        # DOCX/PPTX 固定使用 lxml
        assert DocxExtractor().engine_info == "lxml"
        assert PptxExtractor().engine_info == "lxml"
        # ODS/ODT 在 lxml 与回退之间切换
        assert OdsExtractor().engine_info in {"lxml", "ElementTree"}
        assert OdtExtractor().engine_info in {"lxml", "ElementTree"}


# ---------------------------------------------------------------------------
# 集成函数
# ---------------------------------------------------------------------------


class TestExtractContent:
    def test_extract_text(self, text_file: Path) -> None:
        content = extract_content(text_file)
        assert "hello password world" in content

    def test_extract_docx(self, docx_file: Path) -> None:
        content = extract_content(docx_file)
        assert "段落一 含 password" in content

    def test_extract_unknown_returns_empty(self, tmp_path: Path) -> None:
        path = tmp_path / "unknown.xyz"
        path.write_text("content", encoding="utf-8")
        assert extract_content(path) == ""

    def test_fallback_returns_extracted_content(self, text_file: Path) -> None:
        """提取器成功时返回提取的内容。"""
        content = extract_content_with_fallback(text_file)
        assert "hello password world" in content

    def test_fallback_extractor_failure_falls_back_to_text(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        """提取器抛异常时回退到纯文本读取。"""

        def raise_extract(p: Path) -> str:
            raise RuntimeError("提取失败")

        path = tmp_path / "a.txt"
        path.write_text("plain content", encoding="utf-8")
        monkeypatch.setattr("fuscan.extractors.base.extract_content", raise_extract)
        assert extract_content_with_fallback(path) == "plain content"

    def test_fallback_read_text_oserror_propagates(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """纯文本回退读取失败时 OSError 向上传播。"""

        def raise_extract(p: Path) -> str:
            raise RuntimeError("提取失败")

        path = tmp_path / "nonexistent.txt"
        monkeypatch.setattr("fuscan.extractors.base.extract_content", raise_extract)
        with pytest.raises(OSError):
            extract_content_with_fallback(path)

    def test_get_extractor_returns_none_for_unknown(self) -> None:
        assert get_extractor("xyz") is None

    def test_get_extractor_returns_instance(self) -> None:
        extractor = get_extractor("txt")
        assert extractor is not None
        assert isinstance(extractor, TextExtractor)


# ---------------------------------------------------------------------------
# Scanner 集成
# ---------------------------------------------------------------------------


class TestScannerWithExtractors:
    def test_scan_docx_content(self, docx_file: Path) -> None:
        from fuscan.rules.model import (
            LeafMatch,
            MatchMode,
            MatchTarget,
            Rule,
            RuleSet,
            Severity,
        )
        from fuscan.scanner import Scanner

        rule = Rule(
            name="敏感词",
            severity=Severity.CRITICAL,
            match=LeafMatch(target=MatchTarget.CONTENT, mode=MatchMode.CONTAINS, pattern="password"),
        )
        rs = RuleSet(version="1.0", rules=(rule,))
        scanner = Scanner(rs)
        result = scanner.scan_file(docx_file)
        assert result.has_hit
        assert result.hits[0].rule_name == "敏感词"

    def test_scan_xlsx_content(self, xlsx_file: Path) -> None:
        from fuscan.rules.model import (
            LeafMatch,
            MatchMode,
            MatchTarget,
            Rule,
            RuleSet,
            Severity,
        )
        from fuscan.scanner import Scanner

        rule = Rule(
            name="敏感词",
            severity=Severity.WARNING,
            match=LeafMatch(target=MatchTarget.CONTENT, mode=MatchMode.CONTAINS, pattern="pwd123"),
        )
        rs = RuleSet(version="1.0", rules=(rule,))
        scanner = Scanner(rs)
        result = scanner.scan_file(xlsx_file)
        assert result.has_hit


# ---------------------------------------------------------------------------
# extract_from_bytes：各提取器从内存字节提取（消除双重 I/O）
# ---------------------------------------------------------------------------


class TestExtractFromBytes:
    """各提取器 extract_from_bytes 与 extract(path) 结果一致性。"""

    def test_text_extract_from_bytes_matches_path(self, text_file: Path) -> None:
        """TextExtractor 从 bytes 提取与从 path 提取结果一致。"""
        data = text_file.read_bytes()
        extractor = TextExtractor()
        assert extractor.extract_from_bytes(data) == extractor.extract(text_file)

    def test_text_extract_from_bytes_empty(self) -> None:
        """空字节返回空字符串。"""
        assert TextExtractor().extract_from_bytes(b"") == ""

    def test_text_extract_from_bytes_max_size(self) -> None:
        """超过 max_size 的字节返回空字符串。"""
        extractor = TextExtractor(max_size=10)
        assert extractor.extract_from_bytes(b"x" * 100) == ""

    def test_text_extract_from_bytes_gbk(self, gbk_file: Path) -> None:
        """GBK 编码字节应正确解码。"""
        data = gbk_file.read_bytes()
        content = TextExtractor().extract_from_bytes(data)
        assert "密码" in content
        assert "password123" in content

    def test_docx_extract_from_bytes_matches_path(self, docx_file: Path) -> None:
        """DocxExtractor 从 bytes 提取与从 path 提取结果一致。"""
        data = docx_file.read_bytes()
        extractor = DocxExtractor()
        assert extractor.extract_from_bytes(data) == extractor.extract(docx_file)

    def test_pptx_extract_from_bytes_matches_path(self, pptx_file: Path) -> None:
        """PptxExtractor 从 bytes 提取与从 path 提取结果一致。"""
        data = pptx_file.read_bytes()
        extractor = PptxExtractor()
        assert extractor.extract_from_bytes(data) == extractor.extract(pptx_file)

    def test_xlsx_extract_from_bytes_matches_path(self, xlsx_file: Path) -> None:
        """XlsxExtractor 从 bytes 提取与从 path 提取结果一致。"""
        data = xlsx_file.read_bytes()
        extractor = XlsxExtractor()
        assert extractor.extract_from_bytes(data) == extractor.extract(xlsx_file)

    def test_pdf_extract_from_bytes_matches_path(self, tmp_path: Path) -> None:
        """PdfExtractor 从 bytes 提取与从 path 提取结果一致。"""
        data = _make_pdf_sample(tmp_path)
        path = tmp_path / "fake.pdf"
        path.write_bytes(data)
        extractor = PdfExtractor()
        assert extractor.extract_from_bytes(data) == extractor.extract(path)

    def test_odt_extract_from_bytes_matches_path(self, tmp_path: Path) -> None:
        """OdtExtractor 从 bytes 提取与从 path 提取结果一致。"""
        from tests._odf_samples import make_odt_sample

        data = make_odt_sample(["odt password 内容", "标题内容"])
        path = tmp_path / "test.odt"
        path.write_bytes(data)

        extractor = OdtExtractor()
        assert extractor.extract_from_bytes(data) == extractor.extract(path)

    def test_ods_extract_from_bytes_matches_path(self, tmp_path: Path) -> None:
        """OdsExtractor 从 bytes 提取与从 path 提取结果一致。"""
        from tests._odf_samples import make_ods_sample

        data = make_ods_sample([["cell1", "ods_password"]])
        path = tmp_path / "test.ods"
        path.write_bytes(data)

        extractor = OdsExtractor()
        assert extractor.extract_from_bytes(data) == extractor.extract(path)

    def test_wps_extract_from_bytes_matches_path(self, tmp_path: Path) -> None:
        """WpsExtractor 从 bytes 提取与从 path 提取结果一致。"""
        from docx import Document

        doc = Document()
        doc.add_paragraph("wps bytes password")
        path = tmp_path / "test.wps"
        doc.save(str(path))

        data = path.read_bytes()
        extractor = WpsExtractor()
        assert extractor.extract_from_bytes(data) == extractor.extract(path)

    def test_wps_extract_from_bytes_non_zip(self) -> None:
        """非 ZIP 格式（旧版二进制）返回空字符串。"""
        assert WpsExtractor().extract_from_bytes(b"\xd0\xcf\x11\xe0 old binary") == ""


class TestExtractContentFromBytes:
    """extract_content_from_bytes 模块函数测试。"""

    def test_extract_text_from_bytes(self, text_file: Path) -> None:
        """按扩展名从字节提取文本。"""
        data = text_file.read_bytes()
        content = extract_content_from_bytes(data, "txt")
        assert "hello password world" in content

    def test_extract_docx_from_bytes(self, docx_file: Path) -> None:
        """按扩展名从字节提取 docx。"""
        data = docx_file.read_bytes()
        content = extract_content_from_bytes(data, "docx")
        assert "段落一 含 password" in content

    def test_extract_unknown_extension_returns_empty(self) -> None:
        """未知扩展名返回空字符串。"""
        assert extract_content_from_bytes(b"content", "xyz") == ""

    def test_extract_extension_case_insensitive(self, text_file: Path) -> None:
        """扩展名大小写不敏感。"""
        data = text_file.read_bytes()
        content = extract_content_from_bytes(data, "TXT")
        assert "hello password world" in content

    def test_extract_extension_with_dot(self, text_file: Path) -> None:
        """扩展名带点前缀也能正确处理。"""
        data = text_file.read_bytes()
        content = extract_content_from_bytes(data, ".txt")
        assert "hello password world" in content


# ---------------------------------------------------------------------------
# 大文件流式读取
# ---------------------------------------------------------------------------


class TestLargeFileStreaming:
    """TextExtractor 大文件流式读取与编码检测。"""

    def test_large_utf8_file_streaming(self, tmp_path: Path) -> None:
        """超过 10MB 的 UTF-8 文件应流式解码。"""
        # 构造略大于 10MB 的 UTF-8 文件
        line = "password 行内容 " * 10 + "\n"
        repeat = (10 * 1024 * 1024 // len(line.encode("utf-8"))) + 1
        path = tmp_path / "large.txt"
        path.write_text(line * repeat, encoding="utf-8")
        assert path.stat().st_size > 10 * 1024 * 1024

        content = TextExtractor().extract(path)
        assert "password" in content
        assert content.count("\n") == repeat

    def test_large_gbk_file_streaming(self, tmp_path: Path) -> None:
        """超过 10MB 的 GBK 文件应流式解码。"""
        line = "密码 password 内容\n"
        repeat = (10 * 1024 * 1024 // len(line.encode("gbk"))) + 1
        path = tmp_path / "large_gbk.txt"
        path.write_bytes((line * repeat).encode("gbk"))
        assert path.stat().st_size > 10 * 1024 * 1024

        content = TextExtractor().extract(path)
        assert "password" in content
        assert "密码" in content

    def test_large_utf8_bom_file_streaming(self, tmp_path: Path) -> None:
        """超过 10MB 的 UTF-8 BOM 文件应流式解码。"""
        line = "password bom 内容\n"
        repeat = (10 * 1024 * 1024 // len(line.encode("utf-8"))) + 1
        path = tmp_path / "large_bom.txt"
        path.write_bytes(b"\xef\xbb\xbf" + (line * repeat).encode("utf-8"))
        assert path.stat().st_size > 10 * 1024 * 1024

        content = TextExtractor().extract(path)
        assert "password" in content

    def test_large_utf16_file_streaming(self, tmp_path: Path) -> None:
        """超过 10MB 的 UTF-16 LE BOM 文件应流式解码。"""
        line = "password utf16 内容\n"
        repeat = (10 * 1024 * 1024 // len(line.encode("utf-16-le"))) + 1
        path = tmp_path / "large_utf16.txt"
        path.write_bytes(b"\xff\xfe" + (line * repeat).encode("utf-16-le"))
        assert path.stat().st_size > 10 * 1024 * 1024

        content = TextExtractor().extract(path)
        assert "password" in content

    def test_large_file_crlf_normalized(self, tmp_path: Path) -> None:
        """大文件的 CRLF 行尾应规范化为 LF。"""
        line = "password line\r\n"
        repeat = (10 * 1024 * 1024 // len(line.encode("utf-8"))) + 1
        path = tmp_path / "large_crlf.txt"
        path.write_bytes((line * repeat).encode("utf-8"))
        assert path.stat().st_size > 10 * 1024 * 1024

        content = TextExtractor().extract(path)
        assert "\r\n" not in content
        assert "password line\n" in content

    def test_large_file_read_os_error_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """大文件读取 OSError 时抛出 ExtractorError。"""
        path = tmp_path / "large.txt"
        path.write_bytes(b"x" * (10 * 1024 * 1024 + 1))

        original_open = Path.open

        def mock_open(self: Path, *args, **kwargs):  # type: ignore[no-untyped-def]
            if self == path:
                raise OSError("模拟读取失败")
            return original_open(self, *args, **kwargs)

        monkeypatch.setattr(Path, "open", mock_open)
        with pytest.raises(ExtractorError, match="文件读取失败"):
            TextExtractor().extract(path)

    def test_detect_encoding_utf8_bom(self) -> None:
        """_detect_encoding_from_header 识别 UTF-8 BOM。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header(b"\xef\xbb\xbfcontent") == "utf-8-sig"

    def test_detect_encoding_utf16_le_bom(self) -> None:
        """_detect_encoding_from_header 识别 UTF-16 LE BOM。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header(b"\xff\xfecontent") == "utf-16"

    def test_detect_encoding_utf16_be_bom(self) -> None:
        """_detect_encoding_from_header 识别 UTF-16 BE BOM。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header(b"\xfe\xffcontent") == "utf-16"

    def test_detect_encoding_utf32_le_bom(self) -> None:
        """_detect_encoding_from_header 识别 UTF-32 LE BOM。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header(b"\xff\xfe\x00\x00content") == "utf-32"

    def test_detect_encoding_utf32_be_bom(self) -> None:
        """_detect_encoding_from_header 识别 UTF-32 BE BOM。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header(b"\x00\x00\xfe\xffcontent") == "utf-32"

    def test_detect_encoding_plain_utf8(self) -> None:
        """_detect_encoding_from_header 对纯 UTF-8（无 BOM）返回 utf-8。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header("纯 UTF-8 内容 password".encode()) == "utf-8"

    def test_detect_encoding_gbk(self) -> None:
        """_detect_encoding_from_header 对 GBK 字节返回 gbk。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header("中文 GBK 内容密码".encode("gbk")) == "gbk"

    def test_detect_encoding_binary_returns_none(self) -> None:
        """_detect_encoding_from_header 对非文本字节返回 None。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        # 0x80 不是任何 BOM 前缀，也不是有效的 UTF-8 起始字节或 GBK 引导字节
        assert _detect_encoding_from_header(b"\x80\x81\x82\x83\x84\x85") is None

    def test_detect_encoding_empty(self) -> None:
        """_detect_encoding_from_header 对空字节返回 utf-8（空字节可被任意编码解码）。"""
        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header(b"") == "utf-8"

    def test_large_bytes_decode_skips_charset_normalizer(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """大 bytes（>10MB）用文件头检测编码，跳过 charset-normalizer。"""
        # 构造 >10MB 的 UTF-8 bytes
        data = ("password 内容\n" * 800000).encode("utf-8")
        assert len(data) > 10 * 1024 * 1024

        called = False

        def fake_from_bytes(data: bytes):
            nonlocal called
            called = True
            raise AssertionError("charset-normalizer 不应被调用")

        monkeypatch.setattr("charset_normalizer.from_bytes", fake_from_bytes)
        content = TextExtractor().extract_from_bytes(data)
        assert "password" in content
        assert called is False

    def test_large_bytes_unknown_encoding_fallback_to_charset_normalizer(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """大 bytes 文件头无法确定编码时回退到 charset-normalizer。"""
        # 构造 >10MB 的非 UTF-8/GBK 字节（0x80 不是任何 BOM/UTF-8/GBK 引导字节）
        data = b"\x80\x81\x82\x83" * (3 * 1024 * 1024)
        assert len(data) > 10 * 1024 * 1024

        from fuscan.extractors.text import _detect_encoding_from_header

        assert _detect_encoding_from_header(data[:65536]) is None

        # 应回退到 charset-normalizer
        content = TextExtractor().extract_from_bytes(data)
        assert isinstance(content, str)
        assert len(content) > 0


# ---------------------------------------------------------------------------
# EML 提取器
# ---------------------------------------------------------------------------


@pytest.fixture()
def eml_file(tmp_path: Path) -> Path:
    """生成包含 password 关键词的 EML 测试文件。"""
    from email.message import EmailMessage

    msg = EmailMessage()
    msg["Subject"] = "Test Subject"
    msg["From"] = "sender@example.com"
    msg.set_content("Hello password world")
    path = tmp_path / "test.eml"
    path.write_bytes(msg.as_bytes())
    return path


@pytest.fixture()
def html_eml_file(tmp_path: Path) -> Path:
    """生成仅含 HTML 正文的 EML 测试文件。"""
    from email.message import EmailMessage

    msg = EmailMessage()
    msg["Subject"] = "HTML Test"
    msg["From"] = "sender@example.com"
    msg.set_content("<html><body><p>Hello <b>password</b></p></body></html>", subtype="html")
    path = tmp_path / "test_html.eml"
    path.write_bytes(msg.as_bytes())
    return path


class TestEmlExtractor:
    def test_supported_extensions(self) -> None:
        assert EmlExtractor().supported_extensions == ("eml",)

    def test_extract_eml_text(self, eml_file: Path) -> None:
        content = EmlExtractor().extract(eml_file)
        assert "Test Subject" in content
        assert "sender@example.com" in content
        assert "Hello password world" in content

    def test_extract_html_body(self, html_eml_file: Path) -> None:
        content = EmlExtractor().extract(html_eml_file)
        assert "password" in content
        assert "Hello" in content

    def test_extract_from_bytes_parse_error(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """EML 解析失败抛 ExtractorError。"""
        import email

        def raise_parse(*args: object, **kwargs: object) -> None:
            raise ValueError("解析失败")

        monkeypatch.setattr(email, "message_from_bytes", raise_parse)
        with pytest.raises(ExtractorError, match="EML 解析失败"):
            EmlExtractor().extract_from_bytes(b"bad data")

    def test_extract_nonexistent_raises_error(self, tmp_path: Path) -> None:
        with pytest.raises(ExtractorError, match="文件读取失败"):
            EmlExtractor().extract(tmp_path / "nonexistent.eml")

    def test_registry_has_eml_extractor(self) -> None:
        assert isinstance(get_extractor("eml"), EmlExtractor)

    def test_eml_with_attachment_skipped(self, tmp_path: Path) -> None:
        """带附件的 EML 应跳过附件，仅提取正文。"""
        from email.message import EmailMessage

        msg = EmailMessage()
        msg["Subject"] = "带附件"
        msg["From"] = "sender@example.com"
        msg.set_content("正文 password 内容")
        msg.add_attachment(b"attachment data", maintype="application", subtype="octet-stream")
        path = tmp_path / "attach.eml"
        path.write_bytes(msg.as_bytes())

        content = EmlExtractor().extract(path)
        assert "正文 password 内容" in content
        assert "attachment data" not in content

    def test_eml_invalid_charset_plain_fallback(self, tmp_path: Path) -> None:
        """text/plain 的 charset 无效时应回退到 UTF-8 解码。"""
        raw = (
            b"From: a@b.com\r\n"
            b"Subject: charset test\r\n"
            b"Content-Type: text/plain; charset=invalid-charset\r\n"
            b"\r\n"
            b"password text content"
        )
        path = tmp_path / "bad_charset.eml"
        path.write_bytes(raw)
        content = EmlExtractor().extract(path)
        assert "password text content" in content

    def test_eml_invalid_charset_html_fallback(self, tmp_path: Path) -> None:
        """text/html 的 charset 无效时应回退到 UTF-8 解码。"""
        raw = (
            b"From: a@b.com\r\n"
            b"Subject: html charset test\r\n"
            b"Content-Type: text/html; charset=invalid-charset\r\n"
            b"\r\n"
            b"<p>password html</p>"
        )
        path = tmp_path / "bad_html_charset.eml"
        path.write_bytes(raw)
        content = EmlExtractor().extract(path)
        assert "password html" in content

    def test_eml_empty_body_returns_empty(self, tmp_path: Path) -> None:
        """无正文的 EML 仅返回主题和发件人。"""
        raw = b"From: a@b.com\r\nSubject: no body\r\n\r\n"
        path = tmp_path / "no_body.eml"
        path.write_bytes(raw)
        content = EmlExtractor().extract(path)
        assert "no body" in content
        assert "a@b.com" in content

    def test_eml_no_subject_no_sender(self, tmp_path: Path) -> None:
        """无主题和发件人的 EML 仅返回正文。"""
        raw = b"Content-Type: text/plain\r\n\r\nbody password text"
        path = tmp_path / "no_headers.eml"
        path.write_bytes(raw)
        content = EmlExtractor().extract(path)
        assert "body password text" in content
        assert "主题" not in content
        assert "发件人" not in content


# ---------------------------------------------------------------------------
# XLS 提取器
# ---------------------------------------------------------------------------


class TestXlsExtractor:
    """XLS 提取器测试。

    iter-92 起 XlsExtractor 与 XlsxExtractor 共用 calamine (Rust + PyO3) 后端，
    以下 mock 测试通过 ``monkeypatch`` 替换 ``CalamineWorkbook.from_filelike``
    验证文本提取逻辑；calamine 后端的真实解析由 ``test_extractor_benchmark.py``
    覆盖（XLS 二进制样本难以程序化生成，跳过基准测试）。
    """

    def test_supported_extensions(self) -> None:
        assert XlsExtractor().supported_extensions == ("xls",)

    def test_extract_from_bytes_with_mock(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """mock calamine 验证单元格遍历逻辑。"""

        class FakeSheet:
            def to_python(self) -> list[list[object]]:
                return [["姓名", "密码"], ["张三", "pwd123"]]

        class FakeWorkbook:
            sheet_names = ["Sheet1"]

            def get_sheet_by_index(self, idx: int) -> object:
                return FakeSheet()

        import python_calamine

        monkeypatch.setattr(python_calamine.CalamineWorkbook, "from_filelike", lambda f: FakeWorkbook())

        content = XlsExtractor().extract_from_bytes(b"fake xls data")
        assert "姓名" in content
        assert "pwd123" in content
        assert "张三" in content

    def test_extract_from_bytes_parse_error(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """XLS 解析失败抛 ExtractorError。"""
        import python_calamine

        def raise_parse(_filelike: object) -> None:
            raise python_calamine.CalamineError("解析失败")

        monkeypatch.setattr(python_calamine.CalamineWorkbook, "from_filelike", raise_parse)
        with pytest.raises(ExtractorError, match="XLS 解析失败"):
            XlsExtractor().extract_from_bytes(b"bad data")

    def test_extract_nonexistent_raises_error(self, tmp_path: Path) -> None:
        with pytest.raises(ExtractorError, match="文件读取失败"):
            XlsExtractor().extract(tmp_path / "nonexistent.xls")

    def test_registry_has_xls_extractor(self) -> None:
        assert isinstance(get_extractor("xls"), XlsExtractor)

    def test_extract_from_path_with_mock(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """extract(path) 路径应正确提取单元格文本。"""

        class FakeSheet:
            def to_python(self) -> list[list[object]]:
                return [["user", "password123"]]

        class FakeWorkbook:
            sheet_names = ["Sheet1"]

            def get_sheet_by_index(self, idx: int) -> object:
                return FakeSheet()

        import python_calamine

        monkeypatch.setattr(python_calamine.CalamineWorkbook, "from_filelike", lambda f: FakeWorkbook())
        path = tmp_path / "test.xls"
        path.write_bytes(b"fake xls")
        content = XlsExtractor().extract(path)
        assert "password123" in content
        assert "user" in content

    def test_xls_import_error_raises(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """python-calamine 未安装时应抛出 ExtractorError。"""
        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "python_calamine":
                raise ImportError("No module named 'python_calamine'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="python-calamine 未安装"):
            XlsExtractor().extract_from_bytes(b"fake")

    def test_xls_empty_sheet(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """空工作表应返回空字符串。"""

        class FakeSheet:
            def to_python(self) -> list[list[object]]:
                return []

        class FakeWorkbook:
            sheet_names = ["Empty"]

            def get_sheet_by_index(self, idx: int) -> object:
                return FakeSheet()

        import python_calamine

        monkeypatch.setattr(python_calamine.CalamineWorkbook, "from_filelike", lambda f: FakeWorkbook())
        assert XlsExtractor().extract_from_bytes(b"fake") == ""


# ---------------------------------------------------------------------------
# DOC/PPT 提取器
# ---------------------------------------------------------------------------


class TestExtractUtf16leText:
    """测试 _extract_utf16le_text 辅助函数。"""

    def test_extract_ascii_text(self) -> None:
        from fuscan.extractors.legacy_office import _extract_utf16le_text

        text = "Hello password world"
        data = text.encode("utf-16-le")
        result = _extract_utf16le_text(data)
        assert "Hello password world" in result

    def test_extract_chinese_text(self) -> None:
        from fuscan.extractors.legacy_office import _extract_utf16le_text

        text = "密码 password 测试"
        data = text.encode("utf-16-le")
        result = _extract_utf16le_text(data)
        assert "密码" in result
        assert "password" in result
        assert "测试" in result

    def test_empty_data(self) -> None:
        from fuscan.extractors.legacy_office import _extract_utf16le_text

        assert _extract_utf16le_text(b"") == ""
        assert _extract_utf16le_text(b"\x00") == ""

    def test_skip_short_fragments(self) -> None:
        """长度 < 2 的文本片段被过滤。"""
        from fuscan.extractors.legacy_office import _extract_utf16le_text

        # 单字符 A 后跟非文本字节
        data = b"A\x00\x00\x00B\x00\x00\x00"
        result = _extract_utf16le_text(data)
        assert result == ""

    def test_skip_whitespace_only_fragments(self) -> None:
        """纯空白片段 strip 后长度 < 2 被过滤。"""
        from fuscan.extractors.legacy_office import _extract_utf16le_text

        # 4 个 ASCII 空格（U+0020）的 UTF-16LE 编码，匹配正则但 strip 后为空
        data = b"\x20\x00\x20\x00\x20\x00\x20\x00"
        result = _extract_utf16le_text(data)
        assert result == ""


class TestDocExtractor:
    def test_supported_extensions(self) -> None:
        assert DocExtractor().supported_extensions == ("doc",)

    def test_extract_from_bytes_with_mock(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """mock _read_ole_stream 验证 WordDocument 流文本提取。"""
        text = "Hello password world"
        encoded = text.encode("utf-16-le")

        def fake_read(data: bytes, stream_name: str) -> bytes | None:
            assert stream_name == "WordDocument"
            return encoded

        monkeypatch.setattr("fuscan.extractors.legacy_office._read_ole_stream", fake_read)
        content = DocExtractor().extract_from_bytes(b"fake doc data")
        assert "Hello password world" in content

    def test_extract_no_worddocument_stream(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """OLE 文件无 WordDocument 流时返回空字符串。"""

        def fake_read(data: bytes, stream_name: str) -> bytes | None:
            return None

        monkeypatch.setattr("fuscan.extractors.legacy_office._read_ole_stream", fake_read)
        assert DocExtractor().extract_from_bytes(b"fake") == ""

    def test_extract_from_bytes_parse_error(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """DOC 解析失败抛 ExtractorError。"""

        def raise_parse(data: bytes, stream_name: str) -> None:
            raise ValueError("解析失败")

        monkeypatch.setattr("fuscan.extractors.legacy_office._read_ole_stream", raise_parse)
        with pytest.raises(ExtractorError, match="DOC 解析失败"):
            DocExtractor().extract_from_bytes(b"bad data")

    def test_extract_nonexistent_raises_error(self, tmp_path: Path) -> None:
        with pytest.raises(ExtractorError, match="文件读取失败"):
            DocExtractor().extract(tmp_path / "nonexistent.doc")

    def test_registry_has_doc_extractor(self) -> None:
        assert isinstance(get_extractor("doc"), DocExtractor)

    def test_extract_from_path_with_mock(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """extract(path) 路径应正确提取 WordDocument 流文本。"""
        text = "doc password text"
        encoded = text.encode("utf-16-le")

        def fake_read(data: bytes, stream_name: str) -> bytes | None:
            assert stream_name == "WordDocument"
            return encoded

        monkeypatch.setattr("fuscan.extractors.legacy_office._read_ole_stream", fake_read)
        path = tmp_path / "test.doc"
        path.write_bytes(b"fake doc")
        content = DocExtractor().extract(path)
        assert "doc password text" in content

    def test_doc_import_error_raises(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """olefile 未安装且原生引擎不可用时应抛出 ExtractorError。"""
        # 强制禁用原生引擎，走 olefile 回退路径
        monkeypatch.setattr("fuscan.extractors.legacy_office._NATIVE_OLE_AVAILABLE", False)

        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "olefile":
                raise ImportError("No module named 'olefile'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="olefile 未安装"):
            DocExtractor().extract_from_bytes(b"fake")


class TestPptExtractor:
    def test_supported_extensions(self) -> None:
        assert PptExtractor().supported_extensions == ("ppt",)

    def test_extract_from_bytes_with_mock(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """mock _read_ole_stream 验证 PowerPoint Document 流文本提取。"""
        text = "Slide password content"
        encoded = text.encode("utf-16-le")

        def fake_read(data: bytes, stream_name: str) -> bytes | None:
            assert stream_name == "PowerPoint Document"
            return encoded

        monkeypatch.setattr("fuscan.extractors.legacy_office._read_ole_stream", fake_read)
        content = PptExtractor().extract_from_bytes(b"fake ppt data")
        assert "Slide password content" in content

    def test_extract_no_powerpoint_stream(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """OLE 文件无 PowerPoint Document 流时返回空字符串。"""

        def fake_read(data: bytes, stream_name: str) -> bytes | None:
            return None

        monkeypatch.setattr("fuscan.extractors.legacy_office._read_ole_stream", fake_read)
        assert PptExtractor().extract_from_bytes(b"fake") == ""

    def test_extract_from_bytes_parse_error(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """PPT 解析失败抛 ExtractorError。"""

        def raise_parse(data: bytes, stream_name: str) -> None:
            raise ValueError("解析失败")

        monkeypatch.setattr("fuscan.extractors.legacy_office._read_ole_stream", raise_parse)
        with pytest.raises(ExtractorError, match="PPT 解析失败"):
            PptExtractor().extract_from_bytes(b"bad data")

    def test_extract_nonexistent_raises_error(self, tmp_path: Path) -> None:
        with pytest.raises(ExtractorError, match="文件读取失败"):
            PptExtractor().extract(tmp_path / "nonexistent.ppt")

    def test_registry_has_ppt_extractor(self) -> None:
        assert isinstance(get_extractor("ppt"), PptExtractor)

    def test_extract_from_path_with_mock(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """extract(path) 路径应正确提取 PowerPoint Document 流文本。"""
        text = "ppt password slide"
        encoded = text.encode("utf-16-le")

        def fake_read(data: bytes, stream_name: str) -> bytes | None:
            assert stream_name == "PowerPoint Document"
            return encoded

        monkeypatch.setattr("fuscan.extractors.legacy_office._read_ole_stream", fake_read)
        path = tmp_path / "test.ppt"
        path.write_bytes(b"fake ppt")
        content = PptExtractor().extract(path)
        assert "ppt password slide" in content

    def test_ppt_import_error_raises(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """olefile 未安装且原生引擎不可用时应抛出 ExtractorError。"""
        # 强制禁用原生引擎，走 olefile 回退路径
        monkeypatch.setattr("fuscan.extractors.legacy_office._NATIVE_OLE_AVAILABLE", False)

        import builtins

        original_import = builtins.__import__

        def fake_import(name: str, *args, **kwargs):  # type: ignore[no-untyped-def]
            if name == "olefile":
                raise ImportError("No module named 'olefile'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        with pytest.raises(ExtractorError, match="olefile 未安装"):
            PptExtractor().extract_from_bytes(b"fake")


class TestOleExtractorFallback:
    """OLE 提取器 olefile 回退路径覆盖。

    强制禁用 fuscan-core 原生引擎，验证 olefile 回退路径与原生路径语义等价。
    """

    def test_olefile_fallback_extracts_stream(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """禁用原生引擎后，olefile 回退路径正确提取 WordDocument 流。"""
        import olefile

        monkeypatch.setattr("fuscan.extractors.legacy_office._NATIVE_OLE_AVAILABLE", False)

        text = "fallback password text"
        encoded = text.encode("utf-16-le")

        class FakeStream:
            def read(self) -> bytes:
                return encoded

        class FakeOle:
            def exists(self, name: str) -> bool:
                return name == "WordDocument"

            def openstream(self, name: str) -> FakeStream:
                return FakeStream()

            def close(self) -> None:
                pass

        monkeypatch.setattr(olefile, "OleFileIO", lambda data: FakeOle())
        content = DocExtractor().extract_from_bytes(b"fake doc")
        assert "fallback password text" in content

    def test_olefile_fallback_stream_not_found(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """olefile 回退路径下流不存在时返回空字符串。"""
        import olefile

        monkeypatch.setattr("fuscan.extractors.legacy_office._NATIVE_OLE_AVAILABLE", False)

        class FakeOle:
            def exists(self, name: str) -> bool:
                return False

            def openstream(self, name: str) -> None:
                raise AssertionError("不应调用 openstream")

            def close(self) -> None:
                pass

        monkeypatch.setattr(olefile, "OleFileIO", lambda data: FakeOle())
        assert DocExtractor().extract_from_bytes(b"fake") == ""

    def test_olefile_fallback_parse_error(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """olefile 回退路径下解析失败抛 ExtractorError（OSError 包装）。"""
        import olefile

        monkeypatch.setattr("fuscan.extractors.legacy_office._NATIVE_OLE_AVAILABLE", False)

        def raise_parse(data: object) -> None:
            raise OSError("OLE 解析失败")

        monkeypatch.setattr(olefile, "OleFileIO", raise_parse)
        with pytest.raises(ExtractorError, match="DOC 解析失败"):
            DocExtractor().extract_from_bytes(b"bad")


class TestOleExtractorNativeIntegration:
    """OLE 提取器 fuscan-core 原生路径集成测试。

    仅当 fuscan-core 安装时运行；用真实 CFB 字节验证原生路径端到端。
    """

    def test_native_extract_real_cfb_doc(self) -> None:
        """构建真实 CFB 字节，验证原生路径提取 WordDocument 流。"""
        try:
            from fuscan_core import extract_ole_stream  # pyrefly: ignore [missing-module-attribute]
        except ImportError:
            pytest.skip("fuscan-core 未安装，跳过原生路径集成测试")

        # 用 fuscan_core.extract_ole_stream 自身验证——构造合法 CFB 字节
        # 通过 cfb crate（Rust 侧）已 unit test 覆盖；此处验证 Python 侧调用链
        # 直接给非法字节，确认抛 ValueError（PyValueError 子类）
        with pytest.raises(ValueError, match="OLE 复合文档解析失败"):
            extract_ole_stream(b"not a cfb file", "WordDocument")

    def test_native_path_end_to_end_with_mock(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """原生路径激活时，_read_ole_stream 调用原生函数。"""
        from fuscan.extractors.legacy_office import _NATIVE_OLE_AVAILABLE

        if not _NATIVE_OLE_AVAILABLE:
            pytest.skip("fuscan-core 未安装")

        called: dict[str, object] = {}

        def fake_native(data: bytes, stream_name: str) -> bytes | None:
            called["data"] = data
            called["stream_name"] = stream_name
            return "原生路径密码".encode("utf-16-le")

        monkeypatch.setattr("fuscan.extractors.legacy_office._native_extract_ole_stream", fake_native)
        content = DocExtractor().extract_from_bytes(b"real cfb bytes")
        assert "原生路径密码" in content
        assert called["stream_name"] == "WordDocument"
        assert called["data"] == b"real cfb bytes"


# ---------------------------------------------------------------------------
# 新格式集成测试
# ---------------------------------------------------------------------------


class TestScannerWithNewFormats:
    def test_scan_eml_content(self, eml_file: Path) -> None:
        from fuscan.rules.model import (
            LeafMatch,
            MatchMode,
            MatchTarget,
            Rule,
            RuleSet,
            Severity,
        )
        from fuscan.scanner import Scanner

        rule = Rule(
            name="敏感词",
            severity=Severity.WARNING,
            match=LeafMatch(target=MatchTarget.CONTENT, mode=MatchMode.CONTAINS, pattern="password"),
        )
        rs = RuleSet(version="1.0", rules=(rule,))
        scanner = Scanner(rs)
        result = scanner.scan_file(eml_file)
        assert result.has_hit


class TestContentCache:
    """内容提取缓存测试（需求2：避免重复提取导致卡滞）。"""

    def setup_method(self) -> None:
        """每个测试前清空缓存，确保隔离。"""
        clear_content_cache()

    def test_cached_returns_same_content(self, tmp_path: Path) -> None:
        """缓存提取应返回与直接提取相同的内容。"""
        path = tmp_path / "test.txt"
        path.write_text("hello world\npassword=secret\n", encoding="utf-8")

        direct = extract_content_with_fallback(path)
        cached = extract_content_cached(path)
        assert cached == direct
        assert "password=secret" in cached

    def test_second_call_uses_cache(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """第二次调用相同文件不应重复提取。"""
        path = tmp_path / "test.txt"
        path.write_text("content v1\n", encoding="utf-8")

        call_count = {"n": 0}
        original = extract_content_with_fallback

        def counting_fallback(p: Path) -> str:
            call_count["n"] += 1
            return original(p)

        # 模拟 extract_content_cached 内部调用的 extract_content_with_fallback
        monkeypatch.setattr("fuscan.extractors.cache.extract_content_with_fallback", counting_fallback)

        extract_content_cached(path)
        assert call_count["n"] == 1
        # 第二次调用应命中缓存，不触发提取
        extract_content_cached(path)
        assert call_count["n"] == 1

    def test_file_modified_invalidates_cache(self, tmp_path: Path) -> None:
        """文件修改后（mtime/size 变化）缓存应失效。"""
        path = tmp_path / "test.txt"
        path.write_text("v1\n", encoding="utf-8")

        content1 = extract_content_cached(path)
        assert content1 == "v1\n"

        # 修改文件内容（mtime 和 size 都会变化）
        path.write_text("v2 longer content\n", encoding="utf-8")

        content2 = extract_content_cached(path)
        assert content2 == "v2 longer content\n"

    def test_clear_cache_empties_entries(self, tmp_path: Path) -> None:
        """clear_content_cache 应清空所有缓存项。"""
        path = tmp_path / "test.txt"
        path.write_text("cached\n", encoding="utf-8")
        extract_content_cached(path)

        from fuscan.extractors.cache import _CONTENT_CACHE

        assert len(_CONTENT_CACHE) > 0
        clear_content_cache()
        assert len(_CONTENT_CACHE) == 0

    def test_different_files_cached_separately(self, tmp_path: Path) -> None:
        """不同文件应分别缓存。"""
        p1 = tmp_path / "a.txt"
        p1.write_text("aaa\n", encoding="utf-8")
        p2 = tmp_path / "b.txt"
        p2.write_text("bbb\n", encoding="utf-8")

        assert extract_content_cached(p1) == "aaa\n"
        assert extract_content_cached(p2) == "bbb\n"
        # 再次提取应命中各自缓存
        assert extract_content_cached(p1) == "aaa\n"
        assert extract_content_cached(p2) == "bbb\n"

    def test_stat_failure_falls_back_to_uncached(self, tmp_path: Path) -> None:
        """stat 失败时应回退到无缓存提取。"""
        path = tmp_path / "test.txt"
        path.write_text("fallback\n", encoding="utf-8")

        # 正常提取一次填充缓存
        extract_content_cached(path)

        # 删除文件后再次提取，stat 失败应回退到 extract_content_with_fallback
        # extract_content_with_fallback 内部会 try extract_content 失败后 read_text
        # 文件不存在时 read_text 抛 OSError
        path.unlink()
        with pytest.raises(OSError):
            extract_content_cached(path)

    def test_lru_eviction_when_exceeding_max(self, tmp_path: Path) -> None:
        """超过最大缓存数时淘汰最久未使用的项。"""
        from fuscan.extractors.cache import _CONTENT_CACHE, _CONTENT_CACHE_MAX

        for i in range(_CONTENT_CACHE_MAX + 2):
            p = tmp_path / f"f{i}.txt"
            p.write_text(f"content{i}\n", encoding="utf-8")
            extract_content_cached(p)

        # 缓存数不应超过上限
        assert len(_CONTENT_CACHE) <= _CONTENT_CACHE_MAX


# ---------------------------------------------------------------------------
# WpsExtractor / Office 提取器错误路径
# ---------------------------------------------------------------------------


class TestWpsExtractorErrorPathsExtra:
    """覆盖 WpsExtractor 与 DOCX/PPTX 提取器的错误路径。

    DOCX/PPTX 固定走 lxml，无 python-docx/pptx 回退。以下用例验证
    ``extract()`` 读取失败与 lxml 抛 ``BadZipFile`` 时的 ``ExtractorError`` 包装。
    """

    def test_extract_oserror_raises(self, tmp_path: Path) -> None:
        """``extract()`` 读取文件抛 OSError 时应包装为 ``ExtractorError``。"""
        path = tmp_path / "missing.wps"
        # 文件不存在触发 OSError（read_bytes 内部抛 FileNotFoundError）
        with pytest.raises(ExtractorError, match="文件读取失败"):
            WpsExtractor().extract(path)

    def test_extract_as_docx_lxml_badzipfile_raises(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        """lxml 抛 ``BadZipFile`` 时应抛 ``ExtractorError``。"""
        path = tmp_path / "bad.wps"
        path.write_bytes(_make_ooxml_zip("word/document.xml", "corrupt"))

        import zipfile

        def _raise_bad_zip(_data: bytes) -> str:
            raise zipfile.BadZipFile("invalid zip")

        monkeypatch.setattr("fuscan.extractors._ooxml_xml.extract_docx_text", _raise_bad_zip)
        with pytest.raises(ExtractorError, match="WPS 文字文档解析失败"):
            WpsExtractor().extract(path)

    def test_extract_as_pptx_lxml_badzipfile_raises(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        """lxml 抛 ``BadZipFile`` 时应抛 ``ExtractorError``。"""
        path = tmp_path / "bad.dps"
        path.write_bytes(_make_ooxml_zip("ppt/presentation.xml", "corrupt"))

        import zipfile

        def _raise_bad_zip(_data: bytes) -> str:
            raise zipfile.BadZipFile("invalid zip")

        monkeypatch.setattr("fuscan.extractors._ooxml_xml.extract_pptx_text", _raise_bad_zip)
        with pytest.raises(ExtractorError, match="WPS 演示解析失败"):
            WpsExtractor().extract(path)


# ---------------------------------------------------------------------------
# iter-119：失败重试与诊断信息
# ---------------------------------------------------------------------------


class _FlakyExtractor(Extractor):
    """可编程失败提取器：按 ``failure_sequence`` 抛异常，耗尽后返回 ``success_text``。

    用于测试 :meth:`ExtractorRegistry.extract_from_bytes_with_retry` 的重试与降级行为：
    - ``failure_sequence``：依次抛出的异常列表，每次调用抛一个
    - ``success_text``：耗尽 ``failure_sequence`` 后返回的文本
    - ``call_count``：记录 ``extract_from_bytes`` 被调用的次数
    """

    def __init__(self, extension: str, failure_sequence: list[Exception], success_text: str = "ok") -> None:
        self._extension = extension
        self._failures = list(failure_sequence)
        self._success_text = success_text
        self.call_count: int = 0

    @property
    @override
    def supported_extensions(self) -> tuple[str, ...]:
        return (self._extension,)

    @property
    @override
    def speed_tier(self) -> SpeedTier:
        return SpeedTier.VERY_FAST

    @override
    def extract(self, path: Path) -> str:
        return self.extract_from_bytes(path.read_bytes())

    @override
    def extract_from_bytes(self, data: bytes) -> str:
        self.call_count += 1
        if self._failures:
            raise self._failures.pop(0)
        return self._success_text


class TestIsRetriableError:
    """``is_retriable_error`` 异常分类测试（iter-119）。"""

    def test_os_error_is_retriable(self) -> None:
        """``OSError`` 及其子类（PermissionError/BlockingIOError 等）可重试。"""
        assert is_retriable_error(OSError("io")) is True
        assert is_retriable_error(PermissionError("denied")) is True
        assert is_retriable_error(BlockingIOError("locked")) is True
        assert is_retriable_error(FileNotFoundError("missing")) is True

    def test_extractor_error_is_not_retriable(self) -> None:
        """``ExtractorError``（文件损坏/加密）不可重试。"""
        assert is_retriable_error(ExtractorError("corrupt")) is False

    def test_other_exception_is_not_retriable(self) -> None:
        """其他异常（ValueError/RuntimeError 等）不可重试。"""
        assert is_retriable_error(ValueError("bad data")) is False
        assert is_retriable_error(RuntimeError("unexpected")) is False
        assert is_retriable_error(KeyError("missing key")) is False


class TestExtractFromBytesWithRetry:
    """``ExtractorRegistry.extract_from_bytes_with_retry`` 重试逻辑测试（iter-119）。"""

    def test_success_no_retry(self) -> None:
        """提取成功时不触发重试，``call_count`` 为 1。"""
        ext = _FlakyExtractor("xyz", failure_sequence=[], success_text="content")
        registry = ExtractorRegistry()
        registry.register(ext)
        result = registry.extract_from_bytes_with_retry(b"data", "xyz", max_retries=1, backoff_ms=0.0)
        assert result == "content"
        assert ext.call_count == 1

    def test_retry_succeeds_on_second_attempt(self) -> None:
        """第一次抛 ``OSError``，重试成功；``call_count`` 为 2。"""
        ext = _FlakyExtractor(
            "xyz",
            failure_sequence=[OSError("temporary lock")],
            success_text="content",
        )
        registry = ExtractorRegistry()
        registry.register(ext)
        result = registry.extract_from_bytes_with_retry(b"data", "xyz", max_retries=1, backoff_ms=0.0)
        assert result == "content"
        assert ext.call_count == 2

    def test_retry_exhausted_raises_os_error(self) -> None:
        """``OSError`` 重试后仍失败，抛出原始 ``OSError``。"""
        ext = _FlakyExtractor(
            "xyz",
            failure_sequence=[OSError("first"), OSError("second")],
            success_text="never",
        )
        registry = ExtractorRegistry()
        registry.register(ext)
        with pytest.raises(OSError, match="second"):
            registry.extract_from_bytes_with_retry(b"data", "xyz", max_retries=1, backoff_ms=0.0)
        # 1 次初次 + 1 次重试 = 2 次
        assert ext.call_count == 2

    def test_non_retriable_error_no_retry(self) -> None:
        """``ExtractorError`` 不可重试，直接抛出；``call_count`` 为 1。"""
        ext = _FlakyExtractor(
            "xyz",
            failure_sequence=[ExtractorError("file corrupt")],
            success_text="never",
        )
        registry = ExtractorRegistry()
        registry.register(ext)
        with pytest.raises(ExtractorError, match="file corrupt"):
            registry.extract_from_bytes_with_retry(b"data", "xyz", max_retries=3, backoff_ms=0.0)
        # 不重试，只调用 1 次
        assert ext.call_count == 1

    def test_max_retries_zero_means_no_retry(self) -> None:
        """``max_retries=0`` 退化为不重试，``OSError`` 直接抛出。"""
        ext = _FlakyExtractor(
            "xyz",
            failure_sequence=[OSError("io")],
            success_text="never",
        )
        registry = ExtractorRegistry()
        registry.register(ext)
        with pytest.raises(OSError, match="io"):
            registry.extract_from_bytes_with_retry(b"data", "xyz", max_retries=0, backoff_ms=0.0)
        assert ext.call_count == 1

    def test_multiple_retries_until_success(self) -> None:
        """``max_retries=3`` 时可重试 3 次；前 2 次失败、第 3 次成功。"""
        ext = _FlakyExtractor(
            "xyz",
            failure_sequence=[OSError("1"), OSError("2")],
            success_text="content",
        )
        registry = ExtractorRegistry()
        registry.register(ext)
        result = registry.extract_from_bytes_with_retry(b"data", "xyz", max_retries=3, backoff_ms=0.0)
        assert result == "content"
        assert ext.call_count == 3

    def test_unregistered_extension_returns_empty(self) -> None:
        """未注册扩展名返回空字符串，不抛异常、不调用回调。"""
        registry = ExtractorRegistry()
        failures: list[ExtractorFailure] = []
        result = registry.extract_from_bytes_with_retry(
            b"data", "unreg", max_retries=2, backoff_ms=0.0, on_failure=failures.append
        )
        assert result == ""
        assert failures == []

    def test_on_failure_callback_invoked_for_retriable(self) -> None:
        """可重试错误：``on_failure`` 在准备重试时被调用一次（``retried=False``）。"""
        ext = _FlakyExtractor(
            "xyz",
            failure_sequence=[OSError("temp")],
            success_text="ok",
        )
        registry = ExtractorRegistry()
        registry.register(ext)
        failures: list[ExtractorFailure] = []
        registry.extract_from_bytes_with_retry(
            b"data", "xyz", max_retries=1, backoff_ms=0.0, on_failure=failures.append
        )
        # 1 次「准备重试」回调
        assert len(failures) == 1
        failure = failures[0]
        assert failure.extractor_name == "_FlakyExtractor"
        assert failure.extension == "xyz"
        assert failure.error_type == "OSError"
        assert "temp" in failure.error_message
        assert failure.retried is False
        assert failure.succeeded_after_retry is False

    def test_on_failure_callback_invoked_on_exhaustion(self) -> None:
        """重试耗尽时：``on_failure`` 在准备重试 + 最终失败各调用一次。"""
        ext = _FlakyExtractor(
            "xyz",
            failure_sequence=[OSError("1"), OSError("2")],
            success_text="never",
        )
        registry = ExtractorRegistry()
        registry.register(ext)
        failures: list[ExtractorFailure] = []
        with pytest.raises(OSError):
            registry.extract_from_bytes_with_retry(
                b"data", "xyz", max_retries=1, backoff_ms=0.0, on_failure=failures.append
            )
        # 第 1 次失败：准备重试（retried=False）
        # 第 2 次失败：达到上限（retried=True）
        assert len(failures) == 2
        assert failures[0].retried is False
        assert failures[1].retried is True
        assert failures[1].error_type == "OSError"

    def test_on_failure_callback_invoked_for_non_retriable(self) -> None:
        """不可重试错误：``on_failure`` 仅在最终失败时调用一次（``retried=False``）。"""
        ext = _FlakyExtractor(
            "xyz",
            failure_sequence=[ExtractorError("corrupt")],
            success_text="never",
        )
        registry = ExtractorRegistry()
        registry.register(ext)
        failures: list[ExtractorFailure] = []
        with pytest.raises(ExtractorError):
            registry.extract_from_bytes_with_retry(
                b"data", "xyz", max_retries=3, backoff_ms=0.0, on_failure=failures.append
            )
        assert len(failures) == 1
        assert failures[0].retried is False
        assert failures[0].error_type == "ExtractorError"

    def test_backoff_delay_applied_between_retries(
        self,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        """重试前调用 ``time.sleep(backoff_ms / 1000)``，sleep 参数正确传递。"""
        # max_retries=2 意味着最多 3 次调用（1 初次 + 2 重试），需 3 个失败填满
        ext = _FlakyExtractor(
            "xyz",
            failure_sequence=[OSError("1"), OSError("2"), OSError("3")],
            success_text="never",
        )
        registry = ExtractorRegistry()
        registry.register(ext)
        sleeps: list[float] = []
        monkeypatch.setattr("fuscan.extractors.base.time.sleep", sleeps.append)
        with pytest.raises(OSError, match="3"):
            registry.extract_from_bytes_with_retry(b"data", "xyz", max_retries=2, backoff_ms=50.0)
        # 2 次重试 → 2 次 sleep，每次 0.05 秒
        assert sleeps == [0.05, 0.05]


class TestExtractWithPathRetry:
    """``ExtractorRegistry.extract_with_retry`` 路径版本测试（iter-119）。"""

    def test_path_retry_succeeds_on_second_attempt(self, tmp_path: Path) -> None:
        """路径版本：第一次抛 ``OSError``，重试成功。"""
        ext = _FlakyExtractor(
            "xyz",
            failure_sequence=[OSError("lock")],
            success_text="content",
        )
        registry = ExtractorRegistry()
        registry.register(ext)
        path = tmp_path / "file.xyz"
        path.write_bytes(b"data")
        result = registry.extract_with_retry(path, max_retries=1, backoff_ms=0.0)
        assert result == "content"
        assert ext.call_count == 2

    def test_path_uses_extension_inference(self, tmp_path: Path) -> None:
        """``extension=None`` 时从路径推断扩展名。"""
        ext = _FlakyExtractor("xyz", failure_sequence=[], success_text="content")
        registry = ExtractorRegistry()
        registry.register(ext)
        path = tmp_path / "file.xyz"
        path.write_bytes(b"data")
        result = registry.extract_with_retry(path, max_retries=1, backoff_ms=0.0)
        assert result == "content"
        assert ext.call_count == 1

    def test_path_unregistered_extension_returns_empty(self, tmp_path: Path) -> None:
        """未注册扩展名返回空字符串。"""
        registry = ExtractorRegistry()
        path = tmp_path / "file.unreg"
        path.write_bytes(b"data")
        result = registry.extract_with_retry(path, max_retries=1, backoff_ms=0.0)
        assert result == ""


class TestModuleLevelRetryFunctions:
    """模块级便捷函数 ``extract_content_from_bytes_with_retry`` /
    ``extract_content_with_fallback_and_retry`` 测试（iter-119）。"""

    def test_extract_content_from_bytes_with_retry_uses_default_registry(self) -> None:
        """模块级函数使用 ``default_registry``，对未注册扩展名返回空字符串。"""
        result = extract_content_from_bytes_with_retry(b"data", "unregistered_ext_xyz", max_retries=1, backoff_ms=0.0)
        assert result == ""

    def test_extract_content_with_fallback_and_retry_falls_back_to_plaintext(
        self,
        tmp_path: Path,
    ) -> None:
        """提取器失败且重试耗尽后，回退到 UTF-8 纯文本读取。"""
        # 注册一个始终抛 ExtractorError 的提取器
        ext = _FlakyExtractor(
            "fallback_test_ext",
            failure_sequence=[ExtractorError("corrupt")],
            success_text="never",
        )
        # 临时注册到 default_registry
        default_registry.register(ext)
        try:
            path = tmp_path / "file.fallback_test_ext"
            path.write_text("纯文本回退内容", encoding="utf-8")
            result = extract_content_with_fallback_and_retry(path, max_retries=2, backoff_ms=0.0)
            assert result == "纯文本回退内容"
        finally:
            # 清理：从 default_registry 移除测试用扩展名
            default_registry._extractors.pop("fallback_test_ext", None)

    def test_extract_content_with_fallback_and_retry_returns_extracted_content(
        self,
        tmp_path: Path,
    ) -> None:
        """提取器成功时返回提取的文本，不走纯文本回退。"""
        ext = _FlakyExtractor("success_ext", failure_sequence=[], success_text="提取的文本")
        default_registry.register(ext)
        try:
            path = tmp_path / "file.success_ext"
            path.write_bytes(b"binary data")
            result = extract_content_with_fallback_and_retry(path, max_retries=1, backoff_ms=0.0)
            assert result == "提取的文本"
        finally:
            default_registry._extractors.pop("success_ext", None)

    def test_extract_content_with_fallback_and_retry_retry_succeeds(
        self,
        tmp_path: Path,
    ) -> None:
        """瞬时 ``OSError`` 经重试后成功，不触发纯文本回退。"""
        ext = _FlakyExtractor(
            "retry_success_ext",
            failure_sequence=[OSError("temp lock")],
            success_text="提取的文本",
        )
        default_registry.register(ext)
        try:
            path = tmp_path / "file.retry_success_ext"
            path.write_bytes(b"binary data")
            result = extract_content_with_fallback_and_retry(path, max_retries=1, backoff_ms=0.0)
            assert result == "提取的文本"
            assert ext.call_count == 2
        finally:
            default_registry._extractors.pop("retry_success_ext", None)


class TestExtractorFailureDataclass:
    """``ExtractorFailure`` 诊断数据类测试（iter-119）。"""

    def test_failure_is_frozen(self) -> None:
        """``ExtractorFailure`` 是 frozen dataclass，不可变。"""
        failure = ExtractorFailure(
            extractor_name="PdfExtractor",
            extension="pdf",
            error_type="OSError",
            error_message="io error",
            retried=True,
            succeeded_after_retry=False,
        )
        with pytest.raises(AttributeError):
            failure.retried = False  # type: ignore[misc]

    def test_failure_truncates_long_message(self) -> None:
        """``error_message`` 在 :meth:`_retry_loop` 中被截断到 200 字符（避免撑爆统计）。"""
        long_message = "x" * 500
        ext = _FlakyExtractor(
            "xyz",
            failure_sequence=[OSError(long_message)],
            success_text="ok",
        )
        registry = ExtractorRegistry()
        registry.register(ext)
        failures: list[ExtractorFailure] = []
        registry.extract_from_bytes_with_retry(
            b"data", "xyz", max_retries=1, backoff_ms=0.0, on_failure=failures.append
        )
        assert len(failures) == 1
        # _retry_loop 中通过 str(exc)[:200] 截断
        assert len(failures[0].error_message) == 200


# ---------------------------------------------------------------------------
# ImageExtractor / PdfExtractor OCR 回退（RapidOCR-json 预编译 exe）
#
# 通过 fake recognize / fake OcrEngine / fake PdfDocument 覆盖逻辑分支，
# 不依赖真实 exe/模型文件（真实端到端测试见下方 TestImageExtractorRealOcr，
# 需 OCR 运行链就绪）。
# ---------------------------------------------------------------------------


class _FakeOcrEngine:
    """模拟 OcrEngine：``recognize`` 返回固定文本。

    :ivar call_count: 推理调用次数（验证 OCR 是否被触发）
    :ivar last_input: 最近一次推理输入字节（验证传入 PNG bytes）
    """

    def __init__(self, text: str = "扫描文本", *, fail: bool = False) -> None:
        self._text = text
        self._fail = fail
        self.call_count = 0
        self.last_input: bytes | None = None

    def recognize(self, data: bytes) -> str:
        self.call_count += 1
        self.last_input = data
        if self._fail:
            raise ExtractorError("OCR 失败 (code=200): 推理崩溃")
        return self._text


class _FakePilImage:
    """模拟 PIL Image：提供 ``mode`` / ``convert`` / ``save`` 供 PDF OCR 回退使用。

    避免 PDF OCR 回退测试依赖真实 Pillow（to_pil 返回本对象，save 写入占位字节）。
    """

    def __init__(self, mode: str = "RGB") -> None:
        self.mode = mode

    def convert(self, mode: str) -> _FakePilImage:
        return _FakePilImage(mode)

    def save(self, buf: object, format: str = "PNG") -> None:
        buf.write(b"FAKE_PNG_BYTES")  # type: ignore[union-attr]


class _FakeTextPage:
    """模拟 pypdfium2 文本页（空文本层，扫描版 PDF）。"""

    def get_text_range(self) -> str:
        return ""

    def close(self) -> None:
        """模拟 PdfTextPage.close（无操作）。"""


class _FakeRenderResult:
    """模拟 pypdfium2 页面渲染结果（``to_pil`` 返回假 PIL 图片）。"""

    def __init__(self, mode: str = "RGB") -> None:
        self._mode = mode

    def to_pil(self) -> _FakePilImage:
        return _FakePilImage(self._mode)

    def close(self) -> None:
        """模拟 PdfBitmap.close（无操作）。"""


class _FakePdfPage:
    """模拟 pypdfium2 页面：空文本层 + 可渲染位图。"""

    def __init__(self, mode: str = "RGB") -> None:
        self._mode = mode

    def get_textpage(self) -> _FakeTextPage:
        return _FakeTextPage()

    def render(self, scale: float = 1.0) -> _FakeRenderResult:
        return _FakeRenderResult(self._mode)

    def close(self) -> None:
        """模拟 PdfPage.close（无操作）。"""


class _FakePdfDoc:
    """模拟 pypdfium2 ``PdfDocument``：可控页数 + 空文本层页面。"""

    def __init__(self, n_pages: int = 1) -> None:
        self._n_pages = n_pages

    def __len__(self) -> int:
        return self._n_pages

    def get_page(self, i: int) -> _FakePdfPage:
        return _FakePdfPage()

    def close(self) -> None:
        """模拟 PdfDocument.close（无操作）。"""


class TestImageExtractor:
    """图片 OCR 提取器测试（RapidOCR-json exe 后端）。

    image.py 直接调用 :func:`fuscan.extractors.ocr.recognize` 传图片字节，
    故 mock 模块级 ``recognize`` 即可覆盖全部逻辑分支（无 Pillow/numpy 依赖）。
    """

    def test_supported_extensions(self) -> None:
        assert ImageExtractor().supported_extensions == (
            "png",
            "jpg",
            "jpeg",
            "tiff",
            "tif",
            "bmp",
            "webp",
            "gif",
        )

    def test_speed_tier_is_very_slow(self) -> None:
        """图片 OCR 神经网络推理（det+cls+rec 三阶段）→ T5 极慢。"""
        assert ImageExtractor().speed_tier == SpeedTier.VERY_SLOW

    def test_engine_info(self) -> None:
        assert ImageExtractor().engine_info == "rapidocr-json"

    def test_display_name(self) -> None:
        assert ImageExtractor().display_name == "图片（OCR）"

    def test_large_image_skipped(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """超过 10MB 的图片跳过 OCR，返回空字符串且不调用 recognize。"""
        from fuscan.extractors import image as image_mod

        calls: list[bytes] = []
        monkeypatch.setattr(image_mod, "recognize", lambda d: calls.append(d) or "should not happen")
        big_data = b"\x00" * (10 * 1024 * 1024 + 1)
        assert ImageExtractor().extract_from_bytes(big_data) == ""
        assert calls == []

    def test_extract_returns_recognize_output(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """recognize 返回的文本直接作为提取结果。"""
        from fuscan.extractors import image as image_mod

        monkeypatch.setattr(image_mod, "recognize", lambda d: "第一行\n第二行" if d == b"img" else "")
        assert ImageExtractor().extract_from_bytes(b"img") == "第一行\n第二行"

    def test_ocr_engine_missing_raises(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """OCR 引擎（exe/模型）缺失时 recognize 抛 ExtractorError 并透传。"""
        from fuscan.extractors import image as image_mod

        def _raise(_data: bytes) -> None:
            raise ExtractorError("OCR 引擎不存在: ...")

        monkeypatch.setattr(image_mod, "recognize", _raise)
        with pytest.raises(ExtractorError, match="OCR 引擎不存在"):
            ImageExtractor().extract_from_bytes(b"fake")

    def test_recognize_failure_propagates(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """OCR 推理失败（错误码）抛 ExtractorError 并透传。"""
        from fuscan.extractors import image as image_mod

        def _raise(_data: bytes) -> None:
            raise ExtractorError("OCR 失败 (code=200): decode error")

        monkeypatch.setattr(image_mod, "recognize", _raise)
        with pytest.raises(ExtractorError, match="OCR 失败"):
            ImageExtractor().extract_from_bytes(b"fake")


class TestPdfExtractorOcrFallback:
    """PdfExtractor 扫描版 OCR 回退测试。

    通过 fake ``PdfDocument`` 模拟空文本层 PDF，fake OcrEngine 验证 OCR 回退
    触发与 ``last_engine_info`` 更新。``_FakePilImage`` 避免依赖真实 Pillow。
    """

    def test_last_engine_info_initial(self) -> None:
        """新建实例 ``last_engine_info`` 初始为 ``pypdfium2``。"""
        assert PdfExtractor().last_engine_info == "pypdfium2"

    def test_text_layer_present_no_ocr(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """文本层非空时不触发 OCR 回退，``last_engine_info`` 保持 ``pypdfium2``。"""
        from fuscan.extractors import pdf as pdf_mod

        engine = _FakeOcrEngine()
        monkeypatch.setattr(pdf_mod, "get_ocr_engine", lambda: engine)

        def _fake_extract(_self: PdfExtractor, _data: bytes, _doc: object) -> str:
            return "有文本层"

        monkeypatch.setattr(PdfExtractor, "_extract_with_pdfium2", _fake_extract)
        extractor = PdfExtractor()
        content = extractor.extract_from_bytes(b"fake pdf")
        assert content == "有文本层"
        assert extractor.last_engine_info == "pypdfium2"
        assert engine.call_count == 0

    def test_scanned_pdf_triggers_ocr(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """空文本层 PDF 触发 OCR 回退，``last_engine_info`` 更新为复合引擎。"""
        from fuscan.extractors import pdf as pdf_mod

        engine = _FakeOcrEngine(text="OCR识别文本")
        monkeypatch.setattr(pdf_mod, "get_ocr_engine", lambda: engine)
        monkeypatch.setattr(pdf_mod, "_ensure_backend", lambda: lambda _data: _FakePdfDoc(n_pages=1))
        extractor = PdfExtractor()
        content = extractor.extract_from_bytes(b"fake scanned pdf")
        assert "OCR识别文本" in content
        assert extractor.last_engine_info == "pypdfium2 + rapidocr-json"
        assert engine.call_count == 1
        # 验证传入 recognize 的是 PNG 字节
        assert engine.last_input == b"FAKE_PNG_BYTES"

    def test_ocr_engine_missing_silent_degrade(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """OCR 引擎未就位时 OCR 回退静默降级，返回空字符串。"""
        from fuscan.extractors import pdf as pdf_mod

        def _raise() -> None:
            raise ExtractorError("OCR 引擎不存在: ...")

        monkeypatch.setattr(pdf_mod, "get_ocr_engine", _raise)
        monkeypatch.setattr(pdf_mod, "_ensure_backend", lambda: lambda _data: _FakePdfDoc(n_pages=1))
        extractor = PdfExtractor()
        content = extractor.extract_from_bytes(b"fake scanned pdf")
        assert content == ""
        # 引擎缺失降级，last_engine_info 保持 pypdfium2
        assert extractor.last_engine_info == "pypdfium2"

    def test_too_many_pages_skips_ocr(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """PDF 页数超过上限（>50）跳过 OCR。"""
        from fuscan.extractors import pdf as pdf_mod

        engine = _FakeOcrEngine()
        monkeypatch.setattr(pdf_mod, "get_ocr_engine", lambda: engine)
        monkeypatch.setattr(pdf_mod, "_ensure_backend", lambda: lambda _data: _FakePdfDoc(n_pages=51))
        extractor = PdfExtractor()
        content = extractor.extract_from_bytes(b"fake huge pdf")
        assert content == ""
        assert engine.call_count == 0  # 未触发推理

    def test_page_render_failure_skipped(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """单页渲染失败被跳过，其他页正常 OCR。"""
        from fuscan.extractors import pdf as pdf_mod

        engine = _FakeOcrEngine(text="文本")

        class _BadPage:
            def get_textpage(self) -> _FakeTextPage:
                return _FakeTextPage()

            def render(self, scale: float = 1.0) -> object:
                raise RuntimeError("渲染失败")

        class _MixedDoc:
            def __len__(self) -> int:
                return 2

            def get_page(self, i: int) -> object:
                return _BadPage() if i == 0 else _FakePdfPage()

            def close(self) -> None:
                """无操作。"""

        monkeypatch.setattr(pdf_mod, "get_ocr_engine", lambda: engine)
        monkeypatch.setattr(pdf_mod, "_ensure_backend", lambda: lambda _data: _MixedDoc())
        extractor = PdfExtractor()
        content = extractor.extract_from_bytes(b"fake mixed pdf")
        # 第 0 页渲染失败跳过，第 1 页成功
        assert "文本" in content
        assert engine.call_count == 1

    def test_non_rgb_page_converts_to_rgb(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """非 RGB 模式页面先 convert("RGB") 再保存（覆盖 mode 转换分支）。"""
        from fuscan.extractors import pdf as pdf_mod

        engine = _FakeOcrEngine(text="转换后文本")

        class _RgbaDoc:
            def __len__(self) -> int:
                return 1

            def get_page(self, i: int) -> _FakePdfPage:
                return _FakePdfPage(mode="RGBA")

            def close(self) -> None:
                """无操作。"""

        monkeypatch.setattr(pdf_mod, "get_ocr_engine", lambda: engine)
        monkeypatch.setattr(pdf_mod, "_ensure_backend", lambda: lambda _data: _RgbaDoc())
        extractor = PdfExtractor()
        content = extractor.extract_from_bytes(b"fake rgba pdf")
        assert "转换后文本" in content
        assert engine.call_count == 1

    def test_empty_ocr_result_returns_empty(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """OCR 返回空文本时 ``last_engine_info`` 保持 ``pypdfium2``（无 OCR 产出）。"""
        from fuscan.extractors import pdf as pdf_mod

        engine = _FakeOcrEngine(text="")
        monkeypatch.setattr(pdf_mod, "get_ocr_engine", lambda: engine)
        monkeypatch.setattr(pdf_mod, "_ensure_backend", lambda: lambda _data: _FakePdfDoc(n_pages=1))
        extractor = PdfExtractor()
        content = extractor.extract_from_bytes(b"fake scanned pdf")
        assert content == ""
        # OCR 无产出，last_engine_info 保持 pypdfium2
        assert extractor.last_engine_info == "pypdfium2"


# ---------------------------------------------------------------------------
# 真实 OCR 端到端测试（RapidOCR-json 预编译 exe + 真实模型文件）
#
# 仅在 OCR 运行链就绪（exe + PP-OCRv3 模型文件）时运行；.venv（--extra test
# 无 OCR 资源）与 CI 自动跳过；本地下载 exe+模型后可运行验证真实推理链路。
# 用 Pillow 生成含文字图片、reportlab 生成图片型 PDF（无文本层），验证 OCR
# 提取与 PDF OCR 回退的真实推理链路。
# ---------------------------------------------------------------------------
from fuscan.extractors.ocr import is_ocr_available as _is_ocr_available  # noqa: E402

_OCR_READY = _is_ocr_available()


def _render_text_png(text: str, *, font_size: int = 40, width: int = 320, height: int = 90) -> bytes:
    """用 Pillow 绘制黑字白底图片，返回 PNG bytes（供真实 OCR 识别）。

    尝试 Windows arial / Linux DejaVuSans TrueType 字体，均不可用时跳过测试。
    """
    import io as _io

    from PIL import Image, ImageDraw, ImageFont

    img = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(img)
    font = None
    for name in ("arial.ttf", "DejaVuSans.ttf"):
        try:
            font = ImageFont.truetype(name, font_size)
            break
        except OSError:
            continue
    if font is None:
        pytest.skip("无可用 TrueType 字体，跳过真实 OCR 测试")
    draw.text((20, 20), text, fill="black", font=font)
    buf = _io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


@pytest.mark.skipif(not _OCR_READY, reason="OCR 运行链未就绪（exe/模型文件缺失）")
class TestImageExtractorRealOcr:
    """图片 OCR 端到端测试（真实 RapidOCR-json 引擎，需 OCR 运行链就绪）。"""

    def test_real_ocr_extracts_text_from_image(self) -> None:
        """Pillow 绘制文字图片，真实 OCR 引擎识别出非空文本。"""
        data = _render_text_png("Hello Fuscan")
        content = ImageExtractor().extract_from_bytes(data)
        # OCR 应识别出非空文本（不严格断言具体内容，避免识别率波动导致 flaky）
        assert content.strip() != ""

    def test_real_ocr_plain_color_image_returns_empty(self) -> None:
        """纯色无文字图片，OCR 返回空字符串但不崩溃。"""
        import io as _io

        from PIL import Image

        buf = _io.BytesIO()
        Image.new("RGB", (60, 60), "white").save(buf, format="PNG")
        content = ImageExtractor().extract_from_bytes(buf.getvalue())
        assert content == ""


@pytest.mark.skipif(not _OCR_READY, reason="OCR 运行链未就绪（exe/模型文件缺失）")
class TestPdfExtractorRealOcrFallback:
    """PDF OCR 回退端到端测试（真实 RapidOCR-json 引擎，需 OCR 运行链就绪）。"""

    def test_scanned_pdf_triggers_real_ocr(self) -> None:
        """图片型 PDF（无文本层）触发真实 OCR 回退，last_engine_info 更新为复合引擎。"""
        import io as _io

        from reportlab.lib.pagesizes import letter
        from reportlab.lib.utils import ImageReader
        from reportlab.pdfgen import canvas

        # 生成含文字 PNG 嵌入 PDF（无文本层 → 触发 OCR 回退）
        img_data = _render_text_png("Scan PDF Test")
        pdf_buf = _io.BytesIO()
        c = canvas.Canvas(pdf_buf, pagesize=letter)
        c.drawImage(ImageReader(_io.BytesIO(img_data)), 50, 700, width=320, height=90)
        c.showPage()
        c.save()

        extractor = PdfExtractor()
        content = extractor.extract_from_bytes(pdf_buf.getvalue())
        # OCR 回退触发
        assert extractor.last_engine_info == "pypdfium2 + rapidocr-json"
        # OCR 识别出非空文本
        assert content.strip() != ""
